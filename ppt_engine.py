import json
import re
from dataclasses import dataclass
from typing import Any, Dict, List, Optional
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# --------------------------------------------------------------------------
# Module-level definitions required by llm_service.py and streamlit_app.py
# --------------------------------------------------------------------------
FORBIDDEN_METADATA_TERMS = [
    "internal_notes",
    "prompt_injection",
    "system_prompt",
    "confidential_debug",
]

def enforce_topic_isolation(text: str) -> str:
    """Ensures generated text adheres to safety guardrails."""
    if not text:
        return ""
    clean_text = text
    for term in FORBIDDEN_METADATA_TERMS:
        clean_text = clean_text.replace(term, "[filtered]")
    return clean_text

def normalize_slides(slides_data):
    """Normalizes raw slide inputs into a standardized structure."""
    if not isinstance(slides_data, list):
        return []
    normalized = []
    for slide in slides_data:
        if not isinstance(slide, dict):
            continue
        normalized.append({
            "title": slide.get("title", "Untitled Slide"),
            "subtitle": slide.get("subtitle", ""),
            "layout": slide.get("layout", "cards"),
            "cards": slide.get("cards", slide.get("content", [])),
            "bullets": slide.get("bullets", []),
        })
    return normalized

def parse_slide_json_response(raw_response: str) -> Dict[str, Any]:
    """Wrapper function to parse slide JSON using PPTEngine."""
    engine = PPTEngine()
    return engine.parse_slide_json(raw_response)


@dataclass
class PresentationTheme:
    name: str
    is_dark: bool
    bg_color: RGBColor
    card_bg: RGBColor
    border: RGBColor
    text: RGBColor
    muted: RGBColor
    accent: RGBColor


# Default Themes
DARK_THEME = PresentationTheme(
    name="Dark Cyber",
    is_dark=True,
    bg_color=RGBColor(15, 23, 42),      # #0f172a
    card_bg=RGBColor(30, 41, 59),      # #1e293b
    border=RGBColor(51, 65, 85),       # #334155
    text=RGBColor(248, 250, 252),      # #f8fafc
    muted=RGBColor(148, 163, 184),     # #94a3b8
    accent=RGBColor(249, 115, 22),     # #f97316
)

LIGHT_THEME = PresentationTheme(
    name="Light Minimal",
    is_dark=False,
    bg_color=RGBColor(255, 255, 255),   # #ffffff
    card_bg=RGBColor(241, 245, 249),   # #f1f5f9
    border=RGBColor(203, 213, 225),    # #cbd5e1
    text=RGBColor(15, 23, 42),         # #0f172a
    muted=RGBColor(71, 85, 105),       # #475569
    accent=RGBColor(234, 88, 12),      # #ea580c
)

# Added to resolve the "cannot import name 'THEMES'" error
THEMES = {
    "Dark Cyber": DARK_THEME,
    "Light Minimal": LIGHT_THEME,
}

def build_presentation(slides, theme_name="Dark Cyber", topic="Presentation", progress_callback=None):
    """Convenience wrapper to build a presentation from slide lists using PPTEngine."""
    if progress_callback:
        progress_callback("layout", "Initializing presentation builder...")
    
    theme = THEMES.get(theme_name, DARK_THEME)
    engine = PPTEngine(theme=theme)
    
    payload = {"slides": normalize_slides(slides)}
    prs = engine.generate_presentation(payload)
    
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    file_path = output_dir / "Gamma_Style_Presentation.pptx"
    
    if progress_callback:
        progress_callback("visuals", "Saving presentation file...")
        
    prs.save(str(file_path))
    return str(file_path)


class PPTEngine:
    def __init__(self, theme: Optional[PresentationTheme] = None):
        self.prs = Presentation()
        # Set 16:9 Widescreen aspect ratio
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank_layout = self.prs.slide_layouts[6]
        self.current_theme = theme or DARK_THEME

    def set_theme(self, theme: PresentationTheme):
        self.current_theme = theme

    # --------------------------------------------------------------------------
    # LLM Output Cleaning & Resilient Parsing
    # --------------------------------------------------------------------------
    @staticmethod
    def _strip_llm_noise(raw_text: str) -> str:
        """Strips markdown code fences and extraneous pre/post text from LLM output."""
        cleaned = raw_text.strip()
        if "```" in cleaned:
            match = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", cleaned)
            if match:
                cleaned = match.group(1).strip()
        return cleaned

    @staticmethod
    def _json_candidates(raw_text: str) -> List[str]:
        """Extracts potential valid JSON substrings using slice candidate balancing."""
        cleaned = PPTEngine._strip_llm_noise(raw_text)
        candidates = [cleaned]

        start_idx = cleaned.find("{")
        end_idx = cleaned.rfind("}")
        if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
            candidates.append(cleaned[start_idx : end_idx + 1])

        start_arr = cleaned.find("[")
        end_arr = cleaned.rfind("]")
        if start_arr != -1 and end_arr != -1 and end_arr > start_arr:
            candidates.append(cleaned[start_arr : end_arr + 1])

        return candidates

    def parse_slide_json(self, raw_response: str) -> Dict[str, Any]:
        """Attempts to parse raw LLM response into structured presentation data."""
        for candidate in self._json_candidates(raw_response):
            try:
                data = json.loads(candidate)
                if isinstance(data, dict):
                    return data
                elif isinstance(data, list):
                    return {"slides": data}
            except json.JSONDecodeError:
                continue

        raise ValueError("Failed to extract valid JSON presentation data from response.")

    # --------------------------------------------------------------------------
    # Presentation Generation Logic
    # --------------------------------------------------------------------------
    def generate_presentation(self, data: Dict[str, Any]) -> Presentation:
        """Main entry point to build slides from parsed JSON payload."""
        slides_data = data.get("slides", [])
        if not slides_data:
            slides_data = [data]

        for slide_info in slides_data:
            try:
                self._render_slide(slide_info)
            except Exception as err:
                self._render_resilient_fallback(slide_info, str(err))

        return self.prs

    def _apply_background(self, slide):
        """Fills slide background with current theme's background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.current_theme.bg_color

    def _render_slide(self, slide_info: Dict[str, Any]):
        """Routes slide creation based on layout type."""
        layout_type = slide_info.get("layout", "cards").lower()
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)

        title = slide_info.get("title", "Untitled Slide")
        subtitle = slide_info.get("subtitle", "")

        if layout_type == "title":
            self._render_title_layout(slide, title, subtitle)
        elif layout_type == "bullets":
            bullets = slide_info.get("bullets", [])
            self._render_bullet_layout(slide, title, bullets)
        else:
            cards = slide_info.get("cards", slide_info.get("content", []))
            self._render_card_layout(slide, title, cards)

    def _render_title_layout(self, slide, title: str, subtitle: str):
        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(44)
        p_title.font.bold = True
        p_title.font.color.rgb = self.current_theme.accent
        p_title.alignment = PP_ALIGN.CENTER
        p_title.space_after = Pt(16)

        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = self.current_theme.muted
            p_sub.alignment = PP_ALIGN.CENTER

    def _render_bullet_layout(self, slide, title: str, bullets: List[str]):
        self._add_slide_header(slide, title)

        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
        tf = tx_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.current_theme.text
            p.space_after = Pt(14)

    def _render_card_layout(self, slide, title: str, cards: List[Dict[str, Any]]):
        self._add_slide_header(slide, title)

        if not cards:
            return

        num_cards = min(len(cards), 4)
        margin_x = Inches(1.0)
        top = Inches(2.0)
        card_height = Inches(4.8)
        gap = Inches(0.4)
        total_width = Inches(11.333)

        card_width = (total_width - (gap * (num_cards - 1))) / num_cards

        for i in range(num_cards):
            card_data = cards[i]
            left = margin_x + i * (card_width + gap)

            c_title = card_data.get("title", f"Point {i+1}")
            c_body = card_data.get("body", card_data.get("description", ""))
            c_badge = card_data.get("badge", card_data.get("category", None))

            self._add_card(
                slide=slide,
                left=left,
                top=top,
                width=card_width,
                height=card_height,
                title=c_title,
                body=c_body,
                badge=c_badge,
            )

    def _add_slide_header(self, slide, title: str):
        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.current_theme.accent

    def _add_card(
        self,
        slide,
        left: Inches,
        top: Inches,
        width: Inches,
        height: Inches,
        title: str,
        body: str,
        badge: Optional[str] = None,
        theme: Optional[PresentationTheme] = None,
    ):
        if theme is None:
            theme = self.current_theme

        is_dark = getattr(theme, "is_dark", True)
        bg_color = getattr(theme, "card_bg", RGBColor(30, 41, 59) if is_dark else RGBColor(241, 245, 249))
        border_color = getattr(theme, "border", RGBColor(51, 65, 85) if is_dark else RGBColor(203, 213, 225))
        text_color = getattr(theme, "text", RGBColor(248, 250, 252) if is_dark else RGBColor(15, 23, 42))
        muted_text = getattr(theme, "muted", RGBColor(148, 163, 184) if is_dark else RGBColor(71, 85, 105))
        accent_color = getattr(theme, "accent", RGBColor(249, 115, 22))

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(9)
            p_badge.font.bold = True
            p_badge.font.color.rgb = RGBColor(255, 255, 255) if is_dark else text_color
            p_badge.space_after = Pt(6)
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        p_title.text = title
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = accent_color if is_dark else text_color
        p_title.space_after = Pt(4)

        if body:
            p_body = tf.add_paragraph()
            p_body.text = body
            p_body.font.size = Pt(11)
            p_body.font.color.rgb = muted_text

    def _render_resilient_fallback(self, slide_info: Dict[str, Any], error_msg: str):
        slide = self.prs.slides.add_slide(self.blank_layout)
        self._apply_background(slide)

        title = slide_info.get("title", "Slide Content (Recovered)")
        self._add_slide_header(slide, title)

        tx_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
        tf = tx_box.text_frame
        tf.word_wrap = True

        p_err = tf.paragraphs[0]
        p_err.text = "Notice: Dynamic layout recovered due to render exception."
        p_err.font.size = Pt(12)
        p_err.font.italic = True
        p_err.font.color.rgb = self.current_theme.muted
        p_err.space_after = Pt(12)

        for key, val in slide_info.items():
            if key in ["title", "layout"]:
                continue
            p = tf.add_paragraph()
            p.text = f"• {key.capitalize()}: {val}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.current_theme.text
            p.space_after = Pt(8)
