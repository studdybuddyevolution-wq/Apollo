import datetime
import functools
import gc
import hashlib
import hmac
import html
import io
import json
import os
import random
import re
import requests
import smtplib
import tempfile
import time
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
from email.mime.text import MIMEText
import urllib.parse

import extra_streamlit_components as stx
from google import genai as gemini_client_sdk
from google.genai import types as gemini_types
from groq import Groq
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document as LangchainDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PIL import Image

from docx import Document as DocxDocument
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import streamlit as st

# Import Modular Voice Handler
from voice_handler import render_voice_input, run_tts_synthesis
from vision_handler import render_homework_grading_widget, render_image_question_widget

# Import RAG & Kling AI Video Generator module
from video_generator import render_video_generator_ui

# Import User Settings Page
from settings_app import render_settings_page

# Import Adaptive Socratic Tutor (placement test -> tier-calibrated teaching)
from tutor_engine import render_tutor_mode, init_mastery_state

# Import multi-notebook workspace switcher (per-user, NotebookLM-style)
from notebook_manager import (
    NOTEBOOKS_DIR,
    get_active_notebook,
    list_notebooks,
    render_notebook_switcher,
    save_active_notebook,
)

# Import interactive Plotly chart engine
try:
  from charts import render_dynamic_chart_from_text
except ImportError:

  def render_dynamic_chart_from_text(text):
    pass


# 1. Page Configuration & Title
st.set_page_config(layout="wide", page_title="APOLLO OMNI AI", page_icon="⚡")

# 2. Initialize Cookie Manager for Persistent Auth
cookie_manager = stx.CookieManager()
cookies = cookie_manager.get_all()
if cookies is None:
  cookies = {}

# 3. Explicit Key/Token Initialization
try:
  GROQ_API_KEY = st.secrets.get("GROQ_API_KEY", os.getenv("GROQ_API_KEY", ""))
except Exception:
  GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")

try:
  GEMINI_API_KEY = st.secrets.get("GEMINI_API_KEY", os.getenv("GEMINI_API_KEY", ""))
except Exception:
  GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

try:
  OPENROUTER_API_KEY = st.secrets.get(
      "OPENROUTER_API_KEY", os.getenv("OPENROUTER_API_KEY", "")
  )
except Exception:
  OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")

try:
  TAVILY_API_KEY = st.secrets.get(
      "TAVILY_API_KEY", os.getenv("TAVILY_API_KEY", "")
  )
except Exception:
  TAVILY_API_KEY = os.getenv("TAVILY_API_KEY", "")

try:
  KLING_API_KEY = st.secrets.get(
      "KLING_API_KEY", os.getenv("KLING_API_KEY", "")
  )
except Exception:
  KLING_API_KEY = os.getenv("KLING_API_KEY", "")

# Server-side HMAC Secret for Session Cookie Signing
HMAC_SECRET = st.secrets.get("HMAC_SECRET", "apollo_somaiya_secure_secret_key_2026")


def sign_session_token(email: str) -> str:
  """Signs an email string with HMAC-SHA256 to create a secure session cookie."""
  sig = hmac.new(
      HMAC_SECRET.encode("utf-8"), email.encode("utf-8"), hashlib.sha256
  ).hexdigest()
  return f"{email}:{sig}"


def verify_session_token(token: str) -> bool:
  """Verifies an HMAC-SHA256 signed session cookie."""
  if not token or ":" not in token:
    return False
  parts = token.split(":", 1)
  email, sig = parts[0], parts[1]
  expected = hmac.new(
      HMAC_SECRET.encode("utf-8"), email.encode("utf-8"), hashlib.sha256
  ).hexdigest()
  return hmac.compare_digest(sig, expected)


# 4. Resource Caching Pipelines
try:
  import torch
  torch.set_num_threads(max(1, os.cpu_count() or 4))  # use every core for embedding inference
except Exception:
  pass


@st.cache_resource
def get_embedding_model():
  try:
    import torch as _torch
    _device = "cuda" if _torch.cuda.is_available() else "cpu"
  except Exception:
    _device = "cpu"
  return HuggingFaceEmbeddings(
      model_name="sentence-transformers/all-MiniLM-L6-v2",
      model_kwargs={"device": _device},
      encode_kwargs={"normalize_embeddings": True, "batch_size": 64},
  )


@st.cache_resource
def get_text_splitter():
  return RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)


embedder = get_embedding_model()
text_splitter = get_text_splitter()

# 5. State Management Matrix
if "vector_db" not in st.session_state:
  st.session_state.vector_db = None
if "chat_history" not in st.session_state:
  st.session_state.chat_history = []
if "response_time" not in st.session_state:
  st.session_state.response_time = "0.00"
if "source_reference" not in st.session_state:
  st.session_state.source_reference = (
      "<div class='source-box font-mono'>Awaiting vector alignment...</div>"
  )
if "node_count" not in st.session_state:
  st.session_state.node_count = 0
if "active_studio_tool" not in st.session_state:
  st.session_state.active_studio_tool = "Slide Deck"
if "indexed_sources" not in st.session_state:
  # Readable registry of every source that has been embedded into vector_db,
  # e.g. [{"name": "syllabus.pdf", "kind": "file"}, {"name": "https://...", "kind": "web"}]
  st.session_state.indexed_sources = []
if "dialog_open" not in st.session_state:
  # Whether the NotebookLM-style Studio creation dialog is currently open
  st.session_state.dialog_open = False
if "studio_results" not in st.session_state:
  # Persisted output per Studio tool, e.g. studio_results["Mind Map"] = {...}
  st.session_state.studio_results = {}
if "answer_key_enabled" not in st.session_state:
  st.session_state.answer_key_enabled = True
if "chat_pdf_enabled" not in st.session_state:
  st.session_state.chat_pdf_enabled = True
if "planner_result" not in st.session_state:
  st.session_state.planner_result = ""
if "cross_notebook_results" not in st.session_state:
  st.session_state.cross_notebook_results = []

# Adaptive Socratic Tutor: placement-test scores, tiers, chat state, etc.
init_mastery_state()

# Persistent Signed Cookie Auth State Handling
auth_cookie = cookies.get("apollo_somaiya_session")
if "authenticated" not in st.session_state:
  st.session_state.authenticated = verify_session_token(auth_cookie)

if "otp_sent" not in st.session_state:
  st.session_state.otp_sent = False
if "generated_otp" not in st.session_state:
  st.session_state.generated_otp = None
if "otp_timestamp" not in st.session_state:
  st.session_state.otp_timestamp = 0
if "otp_attempts" not in st.session_state:
  st.session_state.otp_attempts = 0
if "user_email" not in st.session_state:
  st.session_state.user_email = ""

# ── Auto-load persisted user profile ──────────────────────────────────────
_PROFILE_DEFAULTS = {
    "full_name": "",
    "university": "Somaiya University",
    "major": "",
    "learning_style": "Visual & Interactive",
    "detail_level": "Intermediate",
    "default_model": "Qwen 3.6 27B (Groq)",
}
if "user_prefs" not in st.session_state:
  _profile_path = "apollo_user_profile.json"
  if os.path.exists(_profile_path):
    try:
      with open(_profile_path, "r", encoding="utf-8") as _pf:
        _loaded_prefs = json.load(_pf)
      st.session_state.user_prefs = {**_PROFILE_DEFAULTS, **_loaded_prefs}
    except Exception:
      st.session_state.user_prefs = dict(_PROFILE_DEFAULTS)
  else:
    st.session_state.user_prefs = dict(_PROFILE_DEFAULTS)

# Voice TTS output toggle (persisted across reruns)
if "voice_output_enabled" not in st.session_state:
  st.session_state.voice_output_enabled = False

# INITIAL SLIDES DEFAULT STATE (UNIVERSAL WELCOME PAGE)
if "slides_data" not in st.session_state:
  st.session_state.slides_data = [{
      "title": "Welcome to Apollo Omni AI",
      "subtitle": "Cognitive Presentation & RAG Studio",
      "image_keyword": (
          "abstract futuristic orange technology grid network minimalist"
      ),
      "image_prompt": (
          "futuristic orange technology grid over a dark navy workspace,"
          " cinematic lighting, photorealistic, no text"
      ),
      "cards": [
          {
              "heading": "Step 1: Index Knowledge",
              "text": (
                  "Use the sidebar to crawl the web with Tavily or upload"
                  " PDFs/TXT files into the vector database."
              ),
          },
          {
              "heading": "Step 2: Define Presentation",
              "text": (
                  "Type any topic and specific instructions into the"
                  " Presentation Studio controls on the right."
              ),
          },
          {
              "heading": "Step 3: Export & Present",
              "text": (
                  "Generate slides backed by real-time RAG context and export"
                  " directly to a Gamma-style .pptx deck."
              ),
          },
      ],
  }]

# 6. Groq LPU Model Matrix
# Note: Groq deprecates model endpoints periodically. Ensure model IDs match live Groq endpoints.
MODEL_OPTIONS = {
    "Qwen 3.6 27B (Groq)": {
        "provider": "groq",
        "model_id": "qwen/qwen3.6-27b",
        "desc": "High speed 60 RPM, 131K context — primary cognitive engine.",
    },
    "GPT-OSS 120B (Groq)": {
        "provider": "groq",
        "model_id": "openai/gpt-oss-120b",
        "desc": "Massive 120B model with 131K context window.",
    },
    "GPT-OSS 20B (Groq)": {
        "provider": "groq",
        "model_id": "openai/gpt-oss-20b",
        "desc": "Fast 20B model with 131K context window.",
    },
    "Llama 3.3 70B (Groq)": {
        "provider": "groq",
        "model_id": "llama-3.3-70b-versatile",
        "desc": "Flagship Llama 3.3 70B versatile model on Groq LPUs.",
    },
    "Groq Compound Mini": {
        "provider": "groq",
        "model_id": "groq/compound-mini",
        "desc": "Ultra-fast Groq Compound Mini (131K context).",
    },
    # Google Gemini free tier (verified live Aug 2026 after a 404 on the old
    # 2.0 Flash ID -- Google's own deprecation message pointed to 3.6 Flash).
    # Pro-tier models remain paid/heavily gated, so only Flash/Flash-Lite are
    # listed. Check https://ai.google.dev/gemini-api/docs/models if this
    # 404s again -- Google rotates these fast.
    "Gemini 3.6 Flash (Google)": {
        "provider": "gemini",
        "model_id": "gemini-3.6-flash",
        "desc": "Current-gen Google Flash model — free tier, 1M context.",
    },
    "Gemini 2.5 Flash (Google)": {
        "provider": "gemini",
        "model_id": "gemini-2.5-flash",
        "desc": "Prior-gen but stable and still free-tier — strong reasoning, 1M context.",
    },
    "Gemini 2.5 Flash-Lite (Google)": {
        "provider": "gemini",
        "model_id": "gemini-2.5-flash-lite",
        "desc": "Fastest, highest free-tier request quota — best under load.",
    },
}


# ─────────────────────────────────────────────────────────────────────────────
# TAVILY CACHED SEARCH  (ttl=1 hour — avoids repeated API calls for same query)
# ─────────────────────────────────────────────────────────────────────────────

@st.cache_data(ttl=3600, show_spinner=False)
def _cached_tavily_search(query: str, api_key: str) -> dict:
  try:
    from tavily import TavilyClient
    client = TavilyClient(api_key=api_key)
    return client.search(
        query=query,
        search_depth="advanced",
        max_results=4,
        include_answer=True,
    )
  except ImportError:
    response = requests.post(
        "https://api.tavily.com/search",
        json={
            "api_key": api_key,
            "query": query,
            "search_depth": "advanced",
            "max_results": 4,
            "include_answer": True,
        },
        timeout=25,
    )
    if response.status_code == 200:
      return response.json()
    return {}
  except Exception:
    return {}


# ─────────────────────────────────────────────────────────────────────────────
# AGENTIC AUTO-SEARCH INTENT DETECTION
# ─────────────────────────────────────────────────────────────────────────────

_SEARCH_TRIGGERS = {
    "latest", "news", "current", "today", "who is", "what is",
    "2024", "2025", "2026", "recent", "now", "live", "real-time",
    "realtime", "breaking", "update", "trending", "happening",
    "this week", "this month", "this year",
}


def _needs_web_search(query: str) -> bool:
  q = query.lower()
  return any(kw in q for kw in _SEARCH_TRIGGERS)


# 6b. Source Registry & Scoped Retrieval Helpers (NotebookLM-style "Sources")
def register_source(name: str, kind: str = "file"):
  """Adds a human-readable source name to the session registry (deduplicated)."""
  if not name:
    return
  existing_names = {s["name"] for s in st.session_state.indexed_sources}
  if name not in existing_names:
    st.session_state.indexed_sources.append({"name": name, "kind": kind})


def get_source_names() -> list[str]:
  return [s["name"] for s in st.session_state.indexed_sources]


_LOADER_BY_SUFFIX = {
    ".pdf": PyPDFLoader,
    ".txt": lambda p: TextLoader(p, encoding="utf-8"),
    ".docx": Docx2txtLoader,
}


def _parse_uploaded_file(file_name: str, file_bytes: bytes) -> tuple[str, list, str | None]:
  """Writes one uploaded file to a temp path, loads it with the matching
  LangChain loader, and returns (file_name, documents, error_message).
  Designed to be safely run inside a ThreadPoolExecutor for parallel parsing
  of multi-file uploads -- each call gets its own temp file, so there's no
  shared state between threads.
  """
  suffix = os.path.splitext(file_name)[1].lower()
  loader_factory = _LOADER_BY_SUFFIX.get(suffix)
  if loader_factory is None:
    return file_name, [], f"Unsupported file type: {suffix or 'unknown'}"

  path = None
  try:
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
      tmp.write(file_bytes)
      path = tmp.name
    loaded = loader_factory(path).load()
    for d in loaded:
      d.metadata["source"] = file_name  # readable name instead of temp path
    if not loaded:
      return file_name, [], "No extractable text found (scanned image PDF?)."
    return file_name, loaded, None
  except Exception as e:
    return file_name, [], str(e)
  finally:
    if path and os.path.exists(path):
      os.unlink(path)


def get_scoped_context(query: str, selected_sources: list[str] | None, k: int = 6) -> str:
  """
  Retrieves top-k relevant chunks from vector_db, optionally restricted to a
  subset of the indexed sources the user picked in a Studio dialog's
  'Sources' selector. Falls back to the full index when nothing is selected
  or everything is selected.
  """
  if st.session_state.vector_db is None:
    return ""
  try:
    all_names = get_source_names()
    restrict = bool(selected_sources) and 0 < len(selected_sources) < len(all_names)
    fetch_k = k * 4 if restrict else k
    retriever = st.session_state.vector_db.as_retriever(search_kwargs={"k": fetch_k})
    nodes = retriever.invoke(query or "summary")
    if restrict:
      nodes = [n for n in nodes if n.metadata.get("source") in selected_sources]
    nodes = nodes[:k]
    return "\n\n".join(
        f"[{n.metadata.get('source', 'Unknown')}]\n{n.page_content}" for n in nodes
    )
  except Exception:
    return ""


def render_sources_and_topic(
    tool_key: str,
    placeholder: str,
    suggestions: list[str],
    label: str = "What should the topic be?",
    area: bool = False,
):
  """
  Renders the NotebookLM/Gemini-style 'Sources' picker + topic field +
  'Things to try' suggestion chips used inside every Studio creation dialog.
  Returns (topic_text, selected_source_names).
  """
  all_sources = get_source_names()
  n_sources = len(all_sources)
  src_label = f"{n_sources} source{'s' if n_sources != 1 else ''}" if n_sources else "No sources"

  with st.expander(f"📎 Sources — {src_label}", expanded=False):
    if not all_sources:
      st.caption(
          "No materials indexed yet. Add PDFs/TXT files or run a web search from"
          " the sidebar — generation will use general knowledge until then."
      )
      selected = []
    else:
      st.caption("Restrict this generation to specific sources, or leave all checked to use everything indexed.")
      selected = [
          name for name in all_sources
          if st.checkbox(name, value=True, key=f"{tool_key}_src_{name}")
      ]

  st.markdown(
      f"<div style='font-size:12px; font-weight:600; color:#e5e7eb; margin: 10px 0 4px 0;'>{label}</div>",
      unsafe_allow_html=True,
  )
  input_key = f"{tool_key}_input"
  if area:
    topic = st.text_area(label, placeholder=placeholder, key=input_key, label_visibility="collapsed", height=90)
  else:
    topic = st.text_input(label, placeholder=placeholder, key=input_key, label_visibility="collapsed")

  if suggestions:
    st.markdown(
        "<div style='font-size:10px; color:#71717a; margin-top: 10px; text-transform: uppercase; letter-spacing: 0.05em;'>Things to try</div>",
        unsafe_allow_html=True,
    )
    for sug_idx, sug in enumerate(suggestions):
      if st.button(f"• {sug}", key=f"{tool_key}_sugg_{sug_idx}", use_container_width=True):
        st.session_state[input_key] = sug
        st.rerun()

  return st.session_state.get(input_key, ""), selected


# 7. Image Engine via Pollinations
_DEFAULT_IMAGE_PROMPT = (
    "cinematic professional presentation visual, dark navy and orange lighting,"
    " high detail, photorealistic, no text, no watermark"
)
_IMAGE_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) Chrome/120.0",
}


def generate_pollinations_image_url(
    prompt: str, width: int = 1024, height: int = 576
) -> str:
  clean_prompt = (prompt or "").strip() or _DEFAULT_IMAGE_PROMPT
  encoded_prompt = urllib.parse.quote(clean_prompt)
  random_seed = random.randint(1, 999_999)
  return (
      f"https://image.pollinations.ai/prompt/{encoded_prompt}"
      f"?width={width}&height={height}&nologo=true&seed={random_seed}"
  )


def _download_image_bytes(url: str):
  """Safely fetch image bytes. Returns None on timeout or any failure."""
  if not url:
    return None
  try:
    resp = requests.get(url, timeout=5, headers=_IMAGE_HEADERS)
    if resp.status_code == 200 and resp.content and len(resp.content) > 5000:
      return resp.content
  except Exception:
    return None
  return None


@st.cache_data(ttl=3600, show_spinner=False)
def fetch_pollinations_image_bytes(
    prompt: str, width: int = 1024, height: int = 576
):
  """Cached Pollinations download keyed by prompt + dimensions (1h TTL)."""
  url = generate_pollinations_image_url(prompt, width=width, height=height)
  return _download_image_bytes(url)


def _slide_image_prompt(slide_info):
  if not isinstance(slide_info, dict):
    return _DEFAULT_IMAGE_PROMPT
  prompt = (
      slide_info.get("image_prompt")
      or slide_info.get("image_keyword")
      or slide_info.get("title")
      or ""
  )
  prompt = str(prompt).strip()
  return prompt or _DEFAULT_IMAGE_PROMPT


def prefetch_slide_images_parallel(slides_data):
  """Download all slide images concurrently. Failures become None."""
  prompts = [_slide_image_prompt(item) for item in slides_data]
  results = [None] * len(prompts)
  if not prompts:
    return results

  def _load(prompt):
    try:
      return fetch_pollinations_image_bytes(prompt)
    except Exception:
      url = generate_pollinations_image_url(prompt)
      return _download_image_bytes(url)

  workers = min(8, max(1, len(prompts)))
  with ThreadPoolExecutor(max_workers=workers) as pool:
    future_map = {
        pool.submit(_load, prompt): idx for idx, prompt in enumerate(prompts)
    }
    for future in as_completed(future_map):
      idx = future_map[future]
      try:
        results[idx] = future.result()
      except Exception:
        results[idx] = None
  gc.collect()
  return results


def fetch_image_by_keyword(keyword):
  """Backward-compatible helper: returns a temp JPEG path or None."""
  image_bytes = fetch_pollinations_image_bytes(
      keyword or _DEFAULT_IMAGE_PROMPT
  )
  if not image_bytes:
    return None
  try:
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp.write(image_bytes)
    tmp.close()
    return tmp.name
  except Exception:
    return None


# 8. PPTX Builder Engine (Gamma AI Style)
def create_gamma_style_pptx(slides_data):
  prs = Presentation()
  prs.slide_width = Inches(13.333)
  prs.slide_height = Inches(7.5)

  BG_COLOR = RGBColor(15, 23, 42)
  CARD_BG = RGBColor(30, 41, 59)
  CARD_BORDER = RGBColor(51, 65, 85)
  ACCENT_COLOR = RGBColor(249, 115, 22)
  TEXT_PRIMARY = RGBColor(248, 250, 252)
  TEXT_MUTED = RGBColor(148, 163, 184)

  blank_layout = prs.slide_layouts[6]
  slide_images = prefetch_slide_images_parallel(slides_data)

  for index, slide_info in enumerate(slides_data):
    if not isinstance(slide_info, dict):
      continue

    slide = prs.slides.add_slide(blank_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    title_text = slide_info.get("title", f"Slide {index+1}")
    subtitle_text = slide_info.get("subtitle", "")
    cards = slide_info.get("cards", [])
    image_bytes = slide_images[index] if index < len(slide_images) else None

    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.font.name = "Arial"

    if subtitle_text:
      p2 = tf.add_paragraph()
      p2.text = subtitle_text
      p2.font.size = Pt(14)
      p2.font.color.rgb = TEXT_MUTED
      p2.font.name = "Arial"

    has_image = False
    if image_bytes:
      try:
        img_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.8),
            Inches(1.8),
            Inches(3.8),
            Inches(5.0),
        )
        img_card.fill.solid()
        img_card.fill.fore_color.rgb = CARD_BG
        img_card.line.color.rgb = CARD_BORDER

        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            Inches(8.95),
            Inches(1.95),
            width=Inches(3.5),
            height=Inches(4.7),
        )
        has_image = True
      except Exception:
        has_image = False

    content_width = Inches(7.6) if has_image else Inches(11.7)

    if isinstance(cards, list) and len(cards) > 0:
      num_cards = min(len(cards), 4)
      card_height = Inches(4.8 / max(num_cards, 1) - 0.15)
      start_top = Inches(1.8)

      for i in range(num_cards):
        card_item = cards[i]
        top_pos = start_top + i * (card_height + Inches(0.15))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8),
            top_pos,
            content_width,
            card_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER

        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.25)
        tf_card.margin_right = Inches(0.25)
        tf_card.margin_top = Inches(0.12)
        tf_card.margin_bottom = Inches(0.12)

        heading = (
            card_item.get("heading", "") if isinstance(card_item, dict) else ""
        )
        text = (
            card_item.get("text", str(card_item))
            if isinstance(card_item, dict)
            else str(card_item)
        )

        p_head = tf_card.paragraphs[0]
        if heading:
          p_head.text = f"▪ {heading}"
          p_head.font.bold = True
          p_head.font.size = Pt(15)
          p_head.font.color.rgb = ACCENT_COLOR
          p_head.font.name = "Arial"

          p_body = tf_card.add_paragraph()
          p_body.text = text
          p_body.font.size = Pt(12)
          p_body.font.color.rgb = TEXT_PRIMARY
          p_body.font.name = "Arial"
        else:
          p_head.text = f"▪ {text}"
          p_head.font.size = Pt(13)
          p_head.font.color.rgb = TEXT_PRIMARY
          p_head.font.name = "Arial"

  path = "apollo_gamma_presentation.pptx"
  prs.save(path)
  gc.collect()
  return path


# 9. Message-format adapter: OpenAI-style {role, content} -> Gemini's
# {role: "user"/"model", parts:[...]} + a separate system_instruction string.
def _messages_to_gemini(messages):
  system_parts = []
  contents = []
  for m in messages:
    role = m.get("role")
    text = m.get("content", "")
    if role == "system":
      system_parts.append(text)
    elif role == "assistant":
      contents.append({"role": "model", "parts": [{"text": text}]})
    else:
      contents.append({"role": "user", "parts": [{"text": text}]})
  system_instruction = "\n\n".join(system_parts) if system_parts else None
  return system_instruction, contents


_GEMINI_FALLBACKS = ["gemini-3.6-flash", "gemini-2.5-flash", "gemini-2.5-flash-lite"]

# When a response gets cut off mid-way by the per-request output cap, we
# automatically ask the model to keep going from exactly where it stopped,
# up to this many extra rounds, instead of just handing back a truncated
# answer. This is what long asks (a full 80-mark exam paper, a long report)
# actually need -- a bigger fixed max_tokens just moves the same cliff edge
# further out; continuation removes the cliff.
_MAX_CONTINUATIONS = 4
_CONTINUE_INSTRUCTION = (
    "Continue EXACTLY where you left off. Do not repeat any earlier content,"
    " do not restart numbering, and do not add a new introduction or"
    " conclusion -- just keep going from the exact cutoff point as if there"
    " was never a break."
)


def _gemini_finish_reason(candidate) -> str:
  return str(getattr(candidate, "finish_reason", "") or "").upper()


def _generate_gemini_stream(messages, gemini_key, primary_model, max_output_tokens=4096):
  if not gemini_key:
    yield "❌ MISSING CONFIGURATION: Please set a valid 'GEMINI_API_KEY' in Streamlit Secrets."
    return

  client = gemini_client_sdk.Client(api_key=gemini_key.strip())
  models_to_try = list(dict.fromkeys([primary_model] + _GEMINI_FALLBACKS))

  last_exception = None
  for model_id in models_to_try:
    try:
      working_messages = list(messages)
      any_yielded = False
      for round_i in range(_MAX_CONTINUATIONS + 1):
        system_instruction, contents = _messages_to_gemini(working_messages)
        config = gemini_types.GenerateContentConfig(
            system_instruction=system_instruction, temperature=0.3, max_output_tokens=max_output_tokens
        )
        round_text = ""
        finish_reason = ""
        for chunk in client.models.generate_content_stream(model=model_id, contents=contents, config=config):
          if chunk.text:
            round_text += chunk.text
            any_yielded = True
            yield chunk.text
          if getattr(chunk, "candidates", None):
            fr = _gemini_finish_reason(chunk.candidates[0])
            if fr:
              finish_reason = fr

        if "MAX_TOKENS" in finish_reason and round_i < _MAX_CONTINUATIONS and round_text:
          working_messages = working_messages + [
              {"role": "assistant", "content": round_text},
              {"role": "user", "content": _CONTINUE_INSTRUCTION},
          ]
          continue
        break

      if any_yielded:
        return
    except Exception as e:
      last_exception = e
      continue

  yield f"❌ Gemini SDK Failure: {str(last_exception)}"


def _generate_gemini_response(messages, gemini_key, primary_model, max_tokens=1200) -> tuple[str | None, str]:
  if not gemini_key:
    return None, "❌ Missing GEMINI_API_KEY in Streamlit secrets."

  client = gemini_client_sdk.Client(api_key=gemini_key.strip())
  models_to_try = list(dict.fromkeys([primary_model] + _GEMINI_FALLBACKS))

  last_err = ""
  for model_id in models_to_try:
    try:
      working_messages = list(messages)
      full_text = ""
      for round_i in range(_MAX_CONTINUATIONS + 1):
        system_instruction, contents = _messages_to_gemini(working_messages)
        config = gemini_types.GenerateContentConfig(
            system_instruction=system_instruction, temperature=0.3, max_output_tokens=max_tokens
        )
        response = client.models.generate_content(model=model_id, contents=contents, config=config)
        piece = response.text or ""
        full_text += piece
        finish_reason = ""
        if getattr(response, "candidates", None):
          finish_reason = _gemini_finish_reason(response.candidates[0])

        if "MAX_TOKENS" in finish_reason and round_i < _MAX_CONTINUATIONS and piece:
          working_messages = working_messages + [
              {"role": "assistant", "content": piece},
              {"role": "user", "content": _CONTINUE_INSTRUCTION},
          ]
          continue
        break

      if full_text.strip():
        return full_text, f"Success ({model_id})"
    except Exception as ex:
      last_err = str(ex)
      continue

  return None, f"Gemini API Error across models: {last_err}"


# 9a. Provider-Routed LLM Streamer with Automatic Same-Provider Model Fallback
# and Automatic Continuation (keeps going past the per-request token cap
# instead of truncating long documents like exam papers or reports).
def generate_llm_stream(messages, groq_key, selected_model_name, gemini_key="", max_tokens=4096):
  model_cfg = MODEL_OPTIONS.get(selected_model_name, {})
  primary_model = model_cfg.get("model_id", "qwen/qwen3.6-27b")
  provider = model_cfg.get("provider", "groq")

  if provider == "gemini":
    yield from _generate_gemini_stream(messages, gemini_key, primary_model, max_output_tokens=max_tokens)
    return

  if not groq_key or not groq_key.startswith("gsk_"):
    yield (
        "❌ MISSING CONFIGURATION: Please set a valid 'GROQ_API_KEY' starting"
        " with 'gsk_' in Streamlit Secrets."
    )
    return

  client = Groq(api_key=groq_key.strip())

  fallback_list = [
      primary_model,
      "qwen/qwen3.6-27b",
      "openai/gpt-oss-120b",
      "openai/gpt-oss-20b",
      "llama-3.3-70b-versatile",
      "groq/compound-mini",
  ]
  models_to_try = list(dict.fromkeys(fallback_list))

  last_exception = None
  for model_id in models_to_try:
    try:
      working_messages = list(messages)
      any_yielded = False
      for round_i in range(_MAX_CONTINUATIONS + 1):
        stream = client.chat.completions.create(
            model=model_id,
            messages=working_messages,
            temperature=0.3,
            max_tokens=max_tokens,
            stream=True,
        )
        round_text = ""
        finish_reason = None
        for chunk in stream:
          token_text = chunk.choices[0].delta.content or ""
          if token_text:
            round_text += token_text
            any_yielded = True
            yield token_text
          fr = chunk.choices[0].finish_reason
          if fr:
            finish_reason = fr

        if finish_reason == "length" and round_i < _MAX_CONTINUATIONS and round_text:
          working_messages = working_messages + [
              {"role": "assistant", "content": round_text},
              {"role": "user", "content": _CONTINUE_INSTRUCTION},
          ]
          continue
        break

      if any_yielded:
        return
    except Exception as e:
      last_exception = e
      continue

  yield f"❌ Groq SDK Failure: {str(last_exception)}"


# 9b. Provider-Routed Non-Streaming LLM Response Helper with Auto-Continuation
def generate_llm_response(messages, groq_key, selected_model_name, max_tokens=1200, gemini_key="") -> tuple[str | None, str]:
  """
  Non-streaming completion wrapper. Routes to Gemini or Groq based on the
  selected model's provider, with automatic fallback within that same
  provider's free-tier model set, AND automatic continuation if a response
  gets cut off by the per-request token cap (e.g. a long exam paper or
  report). Returns (content_text, status_message).
  """
  model_cfg = MODEL_OPTIONS.get(selected_model_name, {})
  primary_model = model_cfg.get("model_id", "qwen/qwen3.6-27b")
  provider = model_cfg.get("provider", "groq")

  if provider == "gemini":
    return _generate_gemini_response(messages, gemini_key, primary_model, max_tokens)

  if not groq_key or not groq_key.startswith("gsk_"):
    return None, "❌ Missing or invalid GROQ_API_KEY starting with 'gsk_' in Streamlit secrets."

  client = Groq(api_key=groq_key.strip())
  fallback_list = [
      primary_model,
      "qwen/qwen3.6-27b",
      "openai/gpt-oss-120b",
      "openai/gpt-oss-20b",
      "llama-3.3-70b-versatile",
      "groq/compound-mini",
  ]
  models_to_try = list(dict.fromkeys(fallback_list))

  last_err = ""
  for model_id in models_to_try:
    try:
      working_messages = list(messages)
      full_text = ""
      for round_i in range(_MAX_CONTINUATIONS + 1):
        completion = client.chat.completions.create(
            model=model_id,
            messages=working_messages,
            temperature=0.3,
            max_tokens=max_tokens,
        )
        choice = completion.choices[0]
        piece = choice.message.content or ""
        full_text += piece
        finish_reason = getattr(choice, "finish_reason", None)

        if finish_reason == "length" and round_i < _MAX_CONTINUATIONS and piece:
          working_messages = working_messages + [
              {"role": "assistant", "content": piece},
              {"role": "user", "content": _CONTINUE_INSTRUCTION},
          ]
          continue
        break

      if full_text.strip():
        return full_text, f"Success ({model_id})"
    except Exception as ex:
      last_err = str(ex)
      continue

  return None, f"Groq API Error across models: {last_err}"


# 10. Robust JSON Parser & Slide Generator for Groq (RAG-Enabled)
# 9c. Auto "Download as Word Doc" for question-paper-style chat requests
_QUESTION_PAPER_KEYWORDS = [
    "question paper", "exam paper", "test paper", "sample paper",
    "practice paper", "mock test", "worksheet", "quiz paper", "question bank",
]


def _looks_like_question_paper(query_text: str) -> bool:
  if not query_text:
    return False
  low = query_text.lower()
  return any(kw in low for kw in _QUESTION_PAPER_KEYWORDS)


_QP_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_QP_QUESTION_LINE_RE = re.compile(r"^\s*(?:Q\.?\s*)?\(?\d+[\.\)]\s+")


def docx_from_chat_answer(title: str, body_markdown: str) -> str:
  """Converts a question-paper-style chat answer into a printable Word
  document: headings render as headings, **bold** renders as bold, and each
  numbered question gets blank ruled space underneath so a student can
  download it and actually write their answers, on-screen or on paper.
  Returns the saved .docx path.
  """
  doc = DocxDocument()
  doc.styles["Normal"].font.name = "Calibri"
  doc.styles["Normal"].font.size = Pt(11)

  heading = doc.add_heading((title or "Question Paper").strip()[:200], level=1)
  heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

  meta_p = doc.add_paragraph()
  meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
  meta_run = meta_p.add_run(
      "Generated by Apollo Omni AI  •  Name: ______________________  •  Date: ____________"
  )
  meta_run.italic = True
  meta_run.font.size = Pt(9)
  doc.add_paragraph()

  for raw_line in (body_markdown or "").splitlines():
    line = raw_line.rstrip()
    if not line.strip():
      continue

    if line.startswith("### "):
      doc.add_heading(line[4:].strip(), level=3)
      continue
    if line.startswith("## "):
      doc.add_heading(line[3:].strip(), level=2)
      continue
    if line.startswith("# "):
      doc.add_heading(line[2:].strip(), level=1)
      continue

    p = doc.add_paragraph()
    pos = 0
    for m in _QP_BOLD_RE.finditer(line):
      if m.start() > pos:
        p.add_run(line[pos:m.start()])
      bold_run = p.add_run(m.group(1))
      bold_run.bold = True
      pos = m.end()
    if pos < len(line):
      p.add_run(line[pos:])

    # Leave ruled writing space under each numbered question.
    if _QP_QUESTION_LINE_RE.match(line):
      for _ in range(3):
        blank = doc.add_paragraph()
        blank.paragraph_format.space_after = Pt(2)
        blank_run = blank.add_run("_" * 90)
        blank_run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
      doc.add_paragraph()

  out_path = os.path.join(tempfile.gettempdir(), f"apollo_paper_{uuid.uuid4().hex[:8]}.docx")
  doc.save(out_path)
  return out_path


def answer_key_docx_from_question_paper(title: str, answer_markdown: str) -> str:
  doc = DocxDocument()
  doc.styles["Normal"].font.name = "Calibri"
  doc.styles["Normal"].font.size = Pt(10)

  heading = doc.add_heading(f"Answer Key: {(title or 'Question Paper').strip()[:180]}", level=1)
  heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
  meta_p = doc.add_paragraph()
  meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
  meta = meta_p.add_run("Generated by Apollo Omni AI  •  Marking scheme / key points")
  meta.italic = True
  meta.font.size = Pt(9)
  doc.add_paragraph()

  for raw_line in (answer_markdown or "").splitlines():
    line = raw_line.rstrip()
    if not line.strip():
      continue
    if line.startswith("### "):
      doc.add_heading(line[4:].strip(), level=3)
      continue
    if line.startswith("## "):
      doc.add_heading(line[3:].strip(), level=2)
      continue
    if line.startswith("# "):
      doc.add_heading(line[2:].strip(), level=1)
      continue

    p = doc.add_paragraph()
    pos = 0
    for m in _QP_BOLD_RE.finditer(line):
      if m.start() > pos:
        p.add_run(line[pos:m.start()])
      run = p.add_run(m.group(1))
      run.bold = True
      pos = m.end()
    if pos < len(line):
      p.add_run(line[pos:])

  out_path = os.path.join(tempfile.gettempdir(), f"apollo_answer_key_{uuid.uuid4().hex[:8]}.docx")
  doc.save(out_path)
  return out_path


def _plain_text_from_markdown(markdown_text: str) -> str:
  text = re.sub(r"```.*?```", "", markdown_text or "", flags=re.DOTALL)
  text = re.sub(r"`([^`]+)`", r"\1", text)
  text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
  text = re.sub(r"\*(.+?)\*", r"\1", text)
  text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text, flags=re.MULTILINE)
  text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
  return text.strip()


def pdf_from_markdown(title: str, body_markdown: str, ruled_questions: bool = False) -> str | None:
  try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
  except Exception:
    return None

  out_path = os.path.join(tempfile.gettempdir(), f"apollo_pdf_{uuid.uuid4().hex[:8]}.pdf")
  doc = SimpleDocTemplate(
      out_path,
      pagesize=A4,
      rightMargin=0.65 * inch,
      leftMargin=0.65 * inch,
      topMargin=0.65 * inch,
      bottomMargin=0.65 * inch,
  )
  styles = getSampleStyleSheet()
  story = [Paragraph(html.escape((title or "Apollo Export").strip()[:180]), styles["Title"]), Spacer(1, 10)]

  for raw_line in (body_markdown or "").splitlines():
    line = raw_line.strip()
    if not line:
      story.append(Spacer(1, 6))
      continue
    level = 0
    if line.startswith("### "):
      level, line = 3, line[4:].strip()
    elif line.startswith("## "):
      level, line = 2, line[3:].strip()
    elif line.startswith("# "):
      level, line = 1, line[2:].strip()
    line = html.escape(_plain_text_from_markdown(line))
    style = styles["Heading2"] if level else styles["BodyText"]
    story.append(Paragraph(line, style))
    if ruled_questions and _QP_QUESTION_LINE_RE.match(raw_line):
      for _ in range(3):
        story.append(Paragraph("_" * 92, styles["BodyText"]))
      story.append(Spacer(1, 8))
  doc.build(story)
  return out_path


def pdf_from_slides_data(slides_data: list[dict]) -> str | None:
  lines = []
  for idx, slide in enumerate(slides_data or [], start=1):
    lines.append(f"## Slide {idx}: {slide.get('title', 'Untitled')}")
    subtitle = slide.get("subtitle", "")
    if subtitle:
      lines.append(subtitle)
    for card in slide.get("cards", []) or []:
      if isinstance(card, dict):
        lines.append(f"**{card.get('heading', 'Detail')}**: {card.get('text', '')}")
      else:
        lines.append(str(card))
    lines.append("")
  return pdf_from_markdown("Apollo Presentation", "\n".join(lines))


def _citation_label(node, idx: int) -> str:
  source = node.metadata.get("source", "Unknown")
  page = node.metadata.get("page")
  if page is not None:
    try:
      return f"{source}, p.{int(page) + 1}"
    except Exception:
      return f"{source}, p.{page}"
  return source


def build_citation_sources(nodes) -> list[dict]:
  sources = []
  seen = set()
  for idx, node in enumerate(nodes or [], start=1):
    source_id = f"S{idx}"
    label = _citation_label(node, idx)
    snippet = re.sub(r"\s+", " ", node.page_content or "").strip()
    dedupe_key = (label, snippet[:220])
    if dedupe_key in seen:
      continue
    seen.add(dedupe_key)
    sources.append({
        "id": source_id,
        "label": label,
        "source": node.metadata.get("source", "Unknown"),
        "page": node.metadata.get("page"),
        "snippet": snippet[:1200],
    })
  return sources


def citation_context_from_sources(citation_sources: list[dict]) -> str:
  return "\n\n".join(
      f"[{src['id']} | {src['label']}]\n{src['snippet']}"
      for src in citation_sources
  )


def render_cited_markdown(content: str, sources: list[dict] | None = None):
  st.markdown(decorate_citation_links(content, sources))
  render_citation_source_panel(sources)


def decorate_citation_links(content: str, sources: list[dict] | None = None) -> str:
  sources = sources or []
  rendered = content or ""
  for src in sources:
    label = f"[source: {src['id']}]"
    page = src.get("page")
    if page is not None:
      try:
        display = f"source: p.{int(page) + 1}"
      except Exception:
        display = f"source: p.{page}"
    else:
      display = f"source: {src['id']}"
    rendered = rendered.replace(label, f"[{display}](#apollo-source-{src['id'].lower()})")
  return rendered


def render_citation_source_panel(sources: list[dict] | None = None):
  sources = sources or []
  if sources:
    with st.expander("Sources used", expanded=False):
      for src in sources:
        st.markdown(
            f"<a id='apollo-source-{src['id'].lower()}'></a>"
            f"**{src['id']} — {html.escape(src['label'])}**",
            unsafe_allow_html=True,
        )
        st.caption(src["snippet"])


FLASHCARD_STORE = "apollo_flashcards.json"


def _user_key(email: str) -> str:
  return email.strip().lower() or "local"


def load_flashcard_store() -> dict:
  if "flashcard_store" not in st.session_state:
    if os.path.exists(FLASHCARD_STORE):
      try:
        with open(FLASHCARD_STORE, "r", encoding="utf-8") as f:
          st.session_state.flashcard_store = json.load(f)
      except Exception:
        st.session_state.flashcard_store = {}
    else:
      st.session_state.flashcard_store = {}
  return st.session_state.flashcard_store


def save_flashcard_store():
  try:
    with open(FLASHCARD_STORE, "w", encoding="utf-8") as f:
      json.dump(st.session_state.get("flashcard_store", {}), f, indent=2)
  except Exception:
    pass


def parse_flashcards(content: str) -> list[dict]:
  cards = []
  current_q = None
  for raw in (content or "").splitlines():
    line = raw.strip().lstrip("-*0123456789. ")
    if line.lower().startswith("q:"):
      current_q = line[2:].strip()
    elif line.lower().startswith("a:") and current_q:
      cards.append({"question": current_q, "answer": line[2:].strip()})
      current_q = None
  if cards:
    return cards

  pattern = re.compile(r"Q(?:uestion)?\s*\d*[:.)]\s*(.*?)\s*A(?:nswer)?\s*\d*[:.)]\s*(.*?)(?=\n\s*Q(?:uestion)?\s*\d*[:.)]|\Z)", re.IGNORECASE | re.DOTALL)
  for q, a in pattern.findall(content or ""):
    q = re.sub(r"\s+", " ", q).strip()
    a = re.sub(r"\s+", " ", a).strip()
    if q and a:
      cards.append({"question": q, "answer": a})
  return cards


def add_flashcards_for_user(email: str, cards: list[dict], topic: str, sources: list[str]):
  if not cards:
    return
  store = load_flashcard_store()
  user_cards = store.setdefault(_user_key(email), [])
  now = datetime.date.today().isoformat()
  existing = {hashlib.sha1((c.get("question", "") + c.get("answer", "")).encode("utf-8")).hexdigest() for c in user_cards}
  for card in cards:
    fingerprint = hashlib.sha1((card["question"] + card["answer"]).encode("utf-8")).hexdigest()
    if fingerprint in existing:
      continue
    user_cards.append({
        "id": "fc_" + uuid.uuid4().hex[:10],
        "question": card["question"],
        "answer": card["answer"],
        "topic": topic or "General",
        "sources": sources or [],
        "created": now,
        "due": now,
        "interval": 0,
        "ease": 2.5,
        "repetitions": 0,
        "last_reviewed": None,
    })
  save_flashcard_store()


def get_user_flashcards(email: str) -> list[dict]:
  return load_flashcard_store().setdefault(_user_key(email), [])


def due_flashcards(email: str) -> list[dict]:
  today = datetime.date.today().isoformat()
  return [c for c in get_user_flashcards(email) if c.get("due", today) <= today]


def review_flashcard(email: str, card_id: str, quality: int):
  cards = get_user_flashcards(email)
  today = datetime.date.today()
  for card in cards:
    if card["id"] != card_id:
      continue
    ease = float(card.get("ease", 2.5))
    reps = int(card.get("repetitions", 0))
    interval = int(card.get("interval", 0))
    if quality < 3:
      reps = 0
      interval = 1
    else:
      if reps == 0:
        interval = 1
      elif reps == 1:
        interval = 6
      else:
        interval = max(1, round(interval * ease))
      reps += 1
    ease = max(1.3, ease + (0.1 - (5 - quality) * (0.08 + (5 - quality) * 0.02)))
    card.update({
        "ease": round(ease, 2),
        "repetitions": reps,
        "interval": interval,
        "last_reviewed": today.isoformat(),
        "due": (today + datetime.timedelta(days=interval)).isoformat(),
    })
    break
  save_flashcard_store()


def _send_plain_email(target_email: str, subject: str, body: str) -> tuple[bool, str | None]:
  try:
    sender_email = st.secrets.get("EMAIL_SENDER", "")
    sender_pass = st.secrets.get("EMAIL_PASSWORD", "")
    if not sender_email or not sender_pass:
      return False, "Email credentials not configured in Streamlit Secrets."
    msg = MIMEText(body, "plain", "utf-8")
    msg["Subject"] = subject
    msg["From"] = sender_email
    msg["To"] = target_email
    server = smtplib.SMTP("smtp.gmail.com", 587)
    server.starttls()
    server.login(sender_email, sender_pass)
    server.sendmail(sender_email, target_email, msg.as_string())
    server.quit()
    return True, None
  except Exception as e:
    return False, str(e)


def build_weekly_digest(email: str) -> str:
  mastery = st.session_state.get("mastery_profile", {})
  cards = get_user_flashcards(email)
  due_count = len(due_flashcards(email))
  active_nb = get_active_notebook(email)
  lines = [
      "Apollo weekly progress digest",
      "",
      f"Notebook: {active_nb['title'] if active_nb else 'Current workspace'}",
      f"Indexed sources: {len(st.session_state.get('indexed_sources', []))}",
      f"Chat turns this notebook: {len(st.session_state.get('chat_history', []))}",
      f"Flashcards saved: {len(cards)}",
      f"Flashcards due today: {due_count}",
      "",
      "Mastery snapshot:",
  ]
  if mastery:
    for rec in sorted(mastery.values(), key=lambda r: r.get("score", 0))[:8]:
      lines.append(f"- {rec.get('display_name', 'Topic')}: {rec.get('score', 0):.0f}/100 ({rec.get('tier', 'Unrated')})")
  else:
    lines.append("- No tutor mastery checks yet.")
  lines.extend(["", "Keep going. Small reviews compound."])
  return "\n".join(lines)


def search_all_notebooks(user_email: str, query: str, embedder, k: int = 3) -> list[dict]:
  if not query.strip():
    return []
  save_active_notebook(user_email)
  results = []
  active_id = st.session_state.get("active_notebook_id")
  for nb in list_notebooks(user_email):
    nb_id = nb["id"]
    try:
      if nb_id == active_id and st.session_state.get("vector_db") is not None:
        db = st.session_state.vector_db
      else:
        index_dir = os.path.join(NOTEBOOKS_DIR, nb_id, "index")
        if not os.path.exists(index_dir):
          continue
        db = FAISS.load_local(index_dir, embedder, allow_dangerous_deserialization=True)
      for node in db.as_retriever(search_kwargs={"k": k}).invoke(query):
        results.append({
            "notebook": nb.get("title", nb_id),
            "source": _citation_label(node, 0),
            "snippet": re.sub(r"\s+", " ", node.page_content or "").strip()[:900],
        })
    except Exception:
      continue
  return results[:12]


def parse_robust_json(raw_text):
  if not raw_text:
    return None

  clean_text = re.sub(r"<think>.*?</think>", "", raw_text, flags=re.DOTALL)
  if "</think>" in clean_text:
    clean_text = clean_text.split("</think>")[-1]

  start = clean_text.find("{")
  end = clean_text.rfind("}")

  if start != -1 and end != -1 and end > start:
    json_candidate = clean_text[start : end + 1]
    try:
      parsed = json.loads(json_candidate)
      if isinstance(parsed, dict):
        if "slides" in parsed and isinstance(parsed["slides"], list):
          return parsed["slides"]
        for val in parsed.values():
          if isinstance(val, list):
            return val
      elif isinstance(parsed, list):
        return parsed
    except json.JSONDecodeError:
      pass

  return None


def generate_slides_with_groq(
    topic, custom_instructions="", context="", groq_key="", user_prefs=None
):
  groq_key = groq_key.strip() if groq_key else ""
  if not groq_key or not groq_key.startswith("gsk_"):
    return (
        None,
        "Missing active GROQ_API_KEY starting with 'gsk_' in Streamlit Secrets.",
    )

  _prefs = user_prefs or {}
  prefs_line = (
      f"Adapt slides for: learning style = '{_prefs.get('learning_style','General')}', "
      f"depth = '{_prefs.get('detail_level','Intermediate')}'.\n\n"
  )

  prompt = f"""{prefs_line}Create an in-depth 4 to 5 slide presentation outline on the topic: '{topic}'.

For EACH slide, also write a concise, descriptive `image_prompt` (12-20 words) that visually represents that slide's specific topic — include subject, setting, lighting, and style. Do not request text, captions, watermarks, or logos in the image. Keep `image_keyword` as a short 3-6 word fallback.

SPECIFIC USER INSTRUCTIONS / FOCUS POINTS:
{custom_instructions if custom_instructions else "None provided."}

INDEXED KNOWLEDGE BASE CONTEXT:
{context if context else "No extra context provided. Use general knowledge."}

OUTPUT RAW JSON ONLY. Do NOT use markdown backticks or explanations. Start with '{{' and end with '}}'.

SCHEMA REQUIRED:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "subtitle": "Informative Subtitle",
      "image_keyword": "short visual keyword",
      "image_prompt": "concise 12-20 word visual description of this slide topic for an AI image model: subject, setting, lighting, style; no text or logos in the image",
      "cards": [
        {{
          "heading": "Subtopic Heading",
          "text": "Detailed multi-sentence content accurately based on the topic and context provided."
        }},
        {{
          "heading": "Key Insight / Mechanic",
          "text": "Comprehensive analysis of facts or concepts mentioned in the context."
        }}
      ]
    }}
  ]
}}"""

  client = Groq(api_key=groq_key)
  models_to_try = [
      "qwen/qwen3.6-27b",
      "openai/gpt-oss-120b",
      "openai/gpt-oss-20b",
      "llama-3.3-70b-versatile",
      "groq/compound-mini",
  ]

  for model_id in models_to_try:
    try:
      completion = client.chat.completions.create(
          model=model_id,
          messages=[
              {
                  "role": "system",
                  "content": (
                      "You are an expert slide deck generator. You output ONLY"
                      " valid raw JSON strictly derived from the given context"
                      " and user instructions. Every slide must include a"
                      " concise topic-specific image_prompt for AI image"
                      " generation (no text in the image)."
                  ),
              },
              {"role": "user", "content": prompt},
          ],
          temperature=0.2,
          max_tokens=4096,
      )

      raw_text = completion.choices[0].message.content or ""
      parsed_slides = parse_robust_json(raw_text)

      if parsed_slides:
        return parsed_slides, f"Success ({model_id})"

    except Exception:
      continue

  return (
      None,
      "Failed to parse slides JSON.",
  )


# 12. Email Dispatcher Function for Auth
def send_otp_email(target_email, otp_code):
  try:
    sender_email = st.secrets.get("EMAIL_SENDER", "")
    sender_pass = st.secrets.get("EMAIL_PASSWORD", "")
    if not sender_email or not sender_pass:
      return False, "Email credentials missing in Streamlit secrets."

    msg = MIMEText(
        f"Your Apollo Omni AI secure access code is: {otp_code}\n\nIf you did"
        " not request this, please ignore this email."
    )
    msg["Subject"] = "APOLLO OMNI - Access Code"
    msg["From"] = sender_email
    msg["To"] = target_email

    server = smtplib.SMTP("smtp.gmail.com", 587)
    server.starttls()
    server.login(sender_email, sender_pass)
    server.send_message(msg)
    server.quit()
    return True, "Success"
  except Exception as e:
    return False, str(e)


def render_progress_dashboard(user_email: str):
  st.markdown(
      "<h2 style='color:#ff8c00; font-family:\"Inter\",sans-serif; font-weight:700; "
      "font-size:22px; letter-spacing:0.08em; margin-bottom:2px;'>📊 PROGRESS DASHBOARD</h2>"
      "<p style='color:#a1a1aa; font-family:\"JetBrains Mono\",monospace; font-size:12px; "
      "margin-top:0;'>Mastery, sources, flashcards, and weekly momentum in one place.</p>",
      unsafe_allow_html=True,
  )

  mastery = st.session_state.get("mastery_profile", {})
  cards = get_user_flashcards(user_email)
  due = due_flashcards(user_email)
  week_start = datetime.date.today() - datetime.timedelta(days=7)
  reviewed_week = [
      c for c in cards
      if c.get("last_reviewed")
      and datetime.date.fromisoformat(c["last_reviewed"]) >= week_start
  ]

  m1, m2, m3, m4 = st.columns(4)
  m1.metric("Topics tracked", len(mastery))
  m2.metric("Sources indexed", len(st.session_state.get("indexed_sources", [])))
  m3.metric("Cards due today", len(due))
  m4.metric("Reviewed this week", len(reviewed_week))

  left, right = st.columns([3, 2])
  with left:
    st.markdown("**Topic Mastery**")
    if mastery:
      try:
        import plotly.express as px
        rows = [{
            "Topic": rec.get("display_name", key),
            "Score": rec.get("score", 0),
            "Tier": rec.get("tier", "Unrated"),
        } for key, rec in mastery.items()]
        fig = px.bar(rows, x="Topic", y="Score", color="Tier", range_y=[0, 100], template="plotly_dark")
        fig.update_layout(height=360, margin=dict(l=10, r=10, t=20, b=10))
        st.plotly_chart(fig, use_container_width=True)
      except Exception:
        for rec in sorted(mastery.values(), key=lambda r: -r.get("score", 0)):
          st.markdown(f"**{rec.get('display_name', 'Topic')}** — {rec.get('score', 0):.0f}/100")
          st.progress(min(1.0, rec.get("score", 0) / 100))
    else:
      st.info("No mastery data yet. Run a Socratic Tutor placement check to start tracking weak spots.")

  with right:
    st.markdown("**Weak Spots**")
    weak = sorted(mastery.values(), key=lambda r: r.get("score", 0))[:5]
    if weak:
      for rec in weak:
        st.markdown(f"- **{rec.get('display_name', 'Topic')}**: {rec.get('score', 0):.0f}/100 ({rec.get('tier', 'Unrated')})")
    else:
      st.caption("Nothing ranked yet.")

    st.markdown("**Weekly Email Digest**")
    digest = build_weekly_digest(user_email)
    with st.expander("Preview digest", expanded=False):
      st.text(digest)
    if st.button("Email me this week's progress", use_container_width=True, type="primary"):
      ok, err = _send_plain_email(user_email, "Apollo weekly progress digest", digest)
      if ok:
        st.success(f"Digest sent to {user_email}.")
      else:
        st.error(f"Couldn't send digest: {err}")


def render_study_planner(user_email: str):
  st.markdown(
      "<h2 style='color:#ff8c00; font-family:\"Inter\",sans-serif; font-weight:700; "
      "font-size:22px; letter-spacing:0.08em; margin-bottom:2px;'>🗓️ EXAM STUDY PLANNER</h2>"
      "<p style='color:#a1a1aa; font-family:\"JetBrains Mono\",monospace; font-size:12px; "
      "margin-top:0;'>Build a day-by-day revision schedule weighted by your Tutor mastery scores.</p>",
      unsafe_allow_html=True,
  )
  mastery = st.session_state.get("mastery_profile", {})
  default_topics = ", ".join(rec.get("display_name", key) for key, rec in mastery.items()) if mastery else ""

  exam_date = st.date_input(
      "Exam date",
      value=datetime.date.today() + datetime.timedelta(days=14),
      min_value=datetime.date.today(),
  )
  topics_text = st.text_area(
      "Topics to include",
      value=default_topics,
      placeholder="e.g., Trees, Graphs, SQL joins, Normalization",
      height=100,
  )
  hours_per_day = st.slider("Study hours per day", 1, 8, 2)

  if st.button("Generate revision plan", use_container_width=True, type="primary"):
    topics = [t.strip() for t in re.split(r"[,;\n]+", topics_text) if t.strip()]
    if not topics:
      st.warning("Add at least one topic.")
    else:
      days = max(1, (exam_date - datetime.date.today()).days)
      scored = []
      for topic in topics:
        rec = mastery.get(topic.lower(), {})
        score = float(rec.get("score", 40))
        scored.append((topic, max(1, 110 - score)))
      plan_lines = [f"# Apollo Revision Plan: {exam_date.isoformat()}", ""]
      ranked = sorted(scored, key=lambda x: -x[1])
      for offset in range(days):
        day = datetime.date.today() + datetime.timedelta(days=offset)
        topic = ranked[offset % len(ranked)][0]
        review_pool = ", ".join(t for t, _ in ranked[:3])
        plan_lines.extend([
            f"## Day {offset + 1} — {day.strftime('%b %d')}",
            f"- Deep work: {topic} ({hours_per_day} hour{'s' if hours_per_day != 1 else ''})",
            "- Active recall: make or review 10 flashcards",
            f"- Quick check: test the weakest current topics ({review_pool})",
            "",
        ])
      plan_lines.extend([
          f"## Exam Eve — {exam_date.strftime('%b %d')}",
          "- Review mistakes, formulas, definitions, and the highest-weight weak spots.",
          "- Stop heavy new learning early and sleep properly.",
      ])
      st.session_state.planner_result = "\n".join(plan_lines)

  if st.session_state.planner_result:
    st.markdown(st.session_state.planner_result)
    st.download_button(
        "Download study plan (.md)",
        st.session_state.planner_result,
        file_name="apollo_exam_study_plan.md",
        mime="text/markdown",
        use_container_width=True,
    )
    pdf_path = pdf_from_markdown("Apollo Exam Study Plan", st.session_state.planner_result)
    if pdf_path and os.path.exists(pdf_path):
      with open(pdf_path, "rb") as f:
        st.download_button(
            "Download study plan (.pdf)",
            f.read(),
            file_name="apollo_exam_study_plan.pdf",
            mime="application/pdf",
            use_container_width=True,
        )


# 13. CSS Styling & Custom UI Layer
st.markdown(
    """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500;700&display=swap');
    
    :root {
      --primary-orange: #ff8c00;
      --surface-black: #0e0e0e;
      --glass-bg: rgba(20, 20, 20, 0.7);
      --glass-border: rgba(255, 140, 0, 0.15);
      --text-color: #e5e2e1;
    }

    .stApp, [data-testid="stAppViewContainer"], [data-testid="stHeader"] {
        background-color: var(--surface-black) !important;
        color: var(--text-color) !important;
        font-family: 'Inter', sans-serif !important;
    }
    [data-testid="stHeader"] { background: transparent !important; border-bottom: none !important; }
    
    h1, h2, h3, h4, h5, h6, p, span, label, li, small, div { color: var(--text-color); }
    .font-mono { font-family: 'JetBrains Mono', monospace !important; }

    .glass-panel {
      background: var(--glass-bg) !important;
      backdrop-filter: blur(12px) !important;
      border: 1px solid var(--glass-border) !important;
      box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.37) !important;
      transition: transform 0.3s cubic-bezier(0.4, 0, 0.2, 1), border-color 0.3s ease !important;
      border-radius: 6px !important;
      padding: 16px !important;
      margin-bottom: 20px !important;
    }
    .glass-panel:hover {
      border-color: rgba(255, 140, 0, 0.4) !important;
    }
    
    .panel-header {
      font-size: 11px !important;
      font-weight: 700 !important;
      letter-spacing: 0.2em !important;
      color: #a1a1aa !important;
      text-transform: uppercase !important;
      margin-bottom: 16px !important;
      border-bottom: 1px solid rgba(255,255,255,0.05) !important;
      padding-bottom: 12px !important;
      display: flex !important;
      align-items: center !important;
      gap: 8px !important;
    }
    .panel-header::before {
      content: '';
      display: inline-block;
      width: 4px;
      height: 12px;
      background-color: var(--primary-orange);
    }

    div[data-testid="stChatInput"] textarea, div[data-testid="stChatInput"] { 
        background-color: rgba(0,0,0,0.8) !important; 
        border-color: rgba(255, 255, 255, 0.1) !important; 
        color: white !important; 
        font-family: 'JetBrains Mono', monospace !important;
    }
    div[data-baseweb="input"] > div, div[data-baseweb="select"] > div, div[data-baseweb="textarea"] > div {
        background-color: rgba(0,0,0,0.8) !important;
        border: 1px solid var(--glass-border) !important;
    }
    div[data-baseweb="input"] input, div[data-baseweb="textarea"] textarea {
        color: white !important;
        background-color: transparent !important;
    }

    .stButton button {
        background: var(--primary-orange) !important;
        color: #000 !important;
        font-size: 12px !important;
        font-weight: 900 !important;
        text-transform: uppercase !important;
        letter-spacing: 0.2em !important;
        border: none !important;
        border-radius: 4px !important;
        transition: transform 0.2s, background 0.2s !important;
    }
    .stButton button:hover {
        background: #ff9d2e !important;
        transform: scale(1.01);
    }
    
    .stButton button[kind="secondary"] {
        background: transparent !important;
        color: #a1a1aa !important;
        border: 1px solid rgba(255,255,255,0.1) !important;
    }

    div[data-testid="stChatMessage"] {
        font-family: 'JetBrains Mono', monospace !important;
        font-size: 13px !important;
        line-height: 1.6 !important;
        background: transparent !important;
        border: none !important;
        border-radius: 0 !important;
        padding: 16px !important;
        margin-bottom: 12px !important;
    }
    div[data-testid="stChatMessage"]:has(div[aria-label="Chat message from user"]) {
        border-left: 2px solid #a1a1aa !important;
        background: rgba(255, 255, 255, 0.02) !important;
    }
    div[data-testid="stChatMessage"]:has(div[aria-label="Chat message from assistant"]) {
        border-left: 2px solid var(--primary-orange) !important;
        background: rgba(255, 140, 0, 0.05) !important;
    }

    .source-box {
        background: rgba(0,0,0,0.6) !important;
        border: 1px solid rgba(255,255,255,0.05) !important;
        padding: 12px !important;
        border-radius: 2px !important;
        font-family: 'JetBrains Mono', monospace !important;
        font-size: 10px !important;
        color: #a1a1aa !important;
        overflow-x: auto;
        max-height: 350px;
    }

    .omni-header {
        display: flex;
        align-items: center;
        justify-content: space-between;
        padding: 16px 32px;
        background: rgba(0, 0, 0, 0.8);
        backdrop-filter: blur(12px);
        border-bottom: 1px solid rgba(255, 255, 255, 0.05);
        margin-top: -60px;
        margin-bottom: 30px;
        position: sticky;
        top: 0;
        z-index: 50;
    }
    .omni-brand { font-size: 20px; font-weight: 700; letter-spacing: 0.2em; color: white; margin:0; line-height: 1.2;}
    .omni-brand span { color: var(--primary-orange); }
    .omni-subtitle { font-size: 10px; color: #71717a; font-family: 'JetBrains Mono', monospace; letter-spacing: 0.1em; text-transform: uppercase; margin:0;}
</style>
""",
    unsafe_allow_html=True,
)

# Custom Header
st.markdown(
    """
    <div class="omni-header">
        <div style="display: flex; align-items: center; gap: 24px;">
            <div>
                <h1 class="omni-brand">APOLLO <span>OMNI</span></h1>
                <p class="omni-subtitle">Cognitive Study & NotebookLM Studio</p>
            </div>
        </div>
        <div style="display: flex; gap: 40px; align-items: center; font-size: 12px; font-weight: 600; text-transform: uppercase; letter-spacing: 0.1em; color: #a1a1aa;">
            <span style="color: #ff8c00; border-bottom: 1px solid #ff8c00; padding-bottom: 4px;">Console</span>
            <span>Studio</span>
            <span>Cognition</span>
        </div>
    </div>
    """,
    unsafe_allow_html=True,
)

# Spacing Adjustments
st.markdown(
    """
<style>
    .block-container {
        padding-top: 1.5rem !important;
        padding-bottom: 1rem !important;
        max-width: 98% !important;
    }
    .omni-header { margin-top: -10px !important; margin-bottom: 20px !important; }
    [data-testid="stSidebar"] {
        background-color: #0a0a0c !important;
        border-right: 1px solid rgba(255, 140, 0, 0.15) !important;
    }
</style>
""",
    unsafe_allow_html=True,
)


# ================= STUDIO CREATION DIALOGS (NotebookLM-style) =================
# Each dialog mirrors the "Sources -> topic -> Things to try -> Generate" flow
# from the reference screenshot. Generating stores the result in
# st.session_state.studio_results[<tool name>] and closes the dialog; the
# persistent Studio panel then renders whatever is in that slot.

def _dialog_footer_generate_cancel(key_prefix: str) -> bool:
  """Shared Generate/Cancel row. Returns True if Generate was clicked."""
  col_gen, col_cancel = st.columns([2, 1])
  with col_cancel:
    if st.button("Cancel", use_container_width=True, key=f"{key_prefix}_cancel"):
      st.session_state.dialog_open = False
      st.rerun()
  with col_gen:
    return st.button("✨ Generate", type="primary", use_container_width=True, key=f"{key_prefix}_generate")


@st.dialog("🎙️ Audio Overview")
def dialog_audio_overview():
  st.caption("Generate a spoken, podcast-style summary of your sources.")
  topic, sel_sources = render_sources_and_topic(
      "audio",
      placeholder="e.g., Summary of the uploaded lecture notes",
      suggestions=[
          "Summarize the key findings of my sources",
          "Give a 2-host podcast style overview",
          "Focus only on the most exam-relevant points",
      ],
  )
  voice_choice = st.selectbox(
      "Narrator Voice:",
      ["en-US-AriaNeural", "en-US-GuyNeural", "en-US-JennyNeural"],
      key="audio_voice_dialog",
  )

  if _dialog_footer_generate_cancel("audio"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    else:
      with st.spinner("⚡ Synthesizing Audio Overview from your sources..."):
        ctx = get_scoped_context(topic or "summary", sel_sources, k=5)
        prompt = (
            f"Create a concise, highly engaging 2-minute spoken study summary for:"
            f" '{topic or 'the indexed materials'}'.\nContext:\n{ctx}"
        )
        script_text, status = generate_llm_response(
            [{"role": "user", "content": prompt}], GROQ_API_KEY, selected_model,
            max_tokens=600, gemini_key=GEMINI_API_KEY,
        )
        if script_text:
          audio_bytes = run_tts_synthesis(script_text, voice=voice_choice)
          st.session_state.studio_results["Audio Overview"] = {
              "script": script_text,
              "audio": audio_bytes,
              "voice": voice_choice,
              "sources": sel_sources,
          }
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Audio overview failed: {status}")


@st.dialog("💻 Slide Deck")
def dialog_slide_deck():
  st.caption("Generate a Gamma-style presentation from your sources.")
  topic, sel_sources = render_sources_and_topic(
      "slides",
      placeholder="e.g. Quantum Computing or Boeing Planes",
      suggestions=[
          "Turn my sources into a presentation",
          "Focus on financial metrics and key breakthroughs",
          "Create a 6-slide executive summary",
      ],
  )
  custom_prompt_input = st.text_area(
      "Custom Prompt / Specific Points (optional):",
      placeholder="e.g., Focus on architectural comparisons, or a specific section only.",
      key="ppt_custom_prompt_in",
      height=70,
  )

  if _dialog_footer_generate_cancel("slides"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    elif not topic:
      st.warning("Please enter a presentation topic.")
    else:
      with st.spinner("Retrieving indexed blocks & generating presentation via Groq..."):
        ppt_context = get_scoped_context(f"{topic} {custom_prompt_input}".strip(), sel_sources, k=6)
        new_slides, status = generate_slides_with_groq(
            topic=topic,
            custom_instructions=custom_prompt_input,
            context=ppt_context,
            groq_key=GROQ_API_KEY,
            user_prefs=st.session_state.get("user_prefs"),
        )
        if new_slides:
          st.session_state.slides_data = new_slides
          st.session_state.studio_results["Slide Deck"] = {"sources": sel_sources}
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Generation Error: {status}")


@st.dialog("🎬 Video Overview")
def dialog_video_overview():
  st.caption("Generate a short AI video overview from your sources.")
  all_sources = get_source_names()
  sel_sources = all_sources
  if all_sources:
    with st.expander(f"📎 Sources — {len(all_sources)} source{'s' if len(all_sources) != 1 else ''}", expanded=False):
      sel_sources = [n for n in all_sources if st.checkbox(n, value=True, key=f"video_src_{n}")]
  else:
    st.caption("No materials indexed yet — the narrated engine will fall back to general knowledge.")

  render_video_generator_ui(
      groq_key=GROQ_API_KEY,
      kling_key=KLING_API_KEY,
      vector_db=st.session_state.vector_db,
      embedder=embedder,
      user_prefs=st.session_state.get("user_prefs"),
      selected_sources=sel_sources,
  )

  if st.button("Done", use_container_width=True, key="video_done"):
    st.session_state.studio_results["Video Overview"] = {"sources": sel_sources}
    st.session_state.dialog_open = False
    st.rerun()


@st.dialog("🧠 Mind Map")
def dialog_mind_map():
  topic, sel_sources = render_sources_and_topic(
      "mindmap",
      placeholder="e.g. Machine Learning Architecture",
      suggestions=[
          "The mind map must be restricted to a specific source",
          "Focus solely on the key concepts of my materials",
          "Create a mind map to help me study for the exam",
      ],
  )

  if _dialog_footer_generate_cancel("mindmap"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    elif not topic:
      st.warning("Please enter a topic for the mind map.")
    else:
      with st.spinner("Generating mind map breakdown..."):
        ctx = get_scoped_context(topic, sel_sources, k=6)
        prompt = (
            f"Generate a structured hierarchical Mermaid mind map for: '{topic}'."
            f"{' Use ONLY the context below.' if ctx else ''}\nContext:\n{ctx}"
        )
        content, status = generate_llm_response(
            [{"role": "user", "content": prompt}], GROQ_API_KEY, selected_model,
            max_tokens=800, gemini_key=GEMINI_API_KEY,
        )
        if content:
          st.session_state.studio_results["Mind Map"] = {"content": content, "sources": sel_sources, "topic": topic}
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Mind map error: {status}")


@st.dialog("📝 Study Reports")
def dialog_study_reports():
  topic, sel_sources = render_sources_and_topic(
      "reports",
      placeholder="e.g., Executive Summary of Indexed Documents",
      suggestions=[
          "Write an executive summary of my sources",
          "Compare and contrast the key arguments",
          "Draft a study report for the exam syllabus",
      ],
  )

  if _dialog_footer_generate_cancel("reports"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    elif not topic:
      st.warning("Please enter a report focus.")
    else:
      with st.spinner("Compiling structured report..."):
        ctx = get_scoped_context(topic, sel_sources, k=6)
        prompt = f"Write an in-depth, beautifully structured Markdown study report on: '{topic}'. Context:\n{ctx}"
        report_md, status = generate_llm_response(
            [{"role": "user", "content": prompt}], GROQ_API_KEY, selected_model,
            max_tokens=1500, gemini_key=GEMINI_API_KEY,
        )
        if report_md:
          st.session_state.studio_results["Study Reports"] = {"content": report_md, "sources": sel_sources, "topic": topic}
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Report error: {status}")


@st.dialog("📇 Flashcards")
def dialog_flashcards():
  topic, sel_sources = render_sources_and_topic(
      "flashcards",
      placeholder="e.g. Key Definitions & Formulas",
      suggestions=[
          "Make flashcards from my sources",
          "Focus on definitions and formulas",
          "Create flashcards for the hardest concepts",
      ],
  )

  if _dialog_footer_generate_cancel("flashcards"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    elif not topic:
      st.warning("Please enter a flashcard topic.")
    else:
      with st.spinner("Creating flashcard deck..."):
        ctx = get_scoped_context(topic, sel_sources, k=6)
        prompt = (
            f"Create 5 Q&A study flashcards for: '{topic}'. Format every card exactly as Q: ... then A: ..."
            f"{' Base them ONLY on this context:' if ctx else ''}\n{ctx}"
        )
        content, status = generate_llm_response(
            [{"role": "user", "content": prompt}], GROQ_API_KEY, selected_model,
            max_tokens=800, gemini_key=GEMINI_API_KEY,
        )
        if content:
          cards = parse_flashcards(content)
          add_flashcards_for_user(st.session_state.user_email, cards, topic, sel_sources)
          st.session_state.studio_results["Flashcards"] = {
              "content": content,
              "sources": sel_sources,
              "topic": topic,
              "cards_added": len(cards),
          }
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Flashcard error: {status}")


@st.dialog("❓ Practice Quiz")
def dialog_practice_quiz():
  topic, sel_sources = render_sources_and_topic(
      "quiz",
      placeholder="e.g. Exam practice questions",
      suggestions=[
          "Quiz me on my sources",
          "Focus on the most exam-relevant material",
          "Create harder, application-style questions",
      ],
  )

  if _dialog_footer_generate_cancel("quiz"):
    if not GROQ_API_KEY or not GROQ_API_KEY.startswith("gsk_"):
      st.error("❌ Missing or invalid GROQ_API_KEY (must start with 'gsk_').")
    elif not topic:
      st.warning("Please enter a quiz topic.")
    else:
      with st.spinner("Generating multiple-choice quiz..."):
        ctx = get_scoped_context(topic, sel_sources, k=6)
        prompt = (
            f"Generate a 3-question multiple choice quiz with answer explanations for: '{topic}'."
            f"{' Base it ONLY on this context:' if ctx else ''}\n{ctx}"
        )
        content, status = generate_llm_response(
            [{"role": "user", "content": prompt}], GROQ_API_KEY, selected_model,
            max_tokens=1000, gemini_key=GEMINI_API_KEY,
        )
        if content:
          st.session_state.studio_results["Practice Quiz"] = {"content": content, "sources": sel_sources, "topic": topic}
          st.session_state.dialog_open = False
          st.rerun()
        else:
          st.error(f"Quiz error: {status}")


_STUDIO_DIALOGS = {
    "Audio Overview": dialog_audio_overview,
    "Slide Deck": dialog_slide_deck,
    "Video Overview": dialog_video_overview,
    "Mind Map": dialog_mind_map,
    "Study Reports": dialog_study_reports,
    "Flashcards": dialog_flashcards,
    "Practice Quiz": dialog_practice_quiz,
}


# ================= AUTHENTICATION GATEKEEPER =================
if not st.session_state.authenticated:
  st.markdown(
      "<div style='text-align: center; margin-top: 80px;'><h2"
      " style='color: #ff8c00; font-family: \"Inter\"; letter-spacing: 0.1em;'>🔒 SECURE ACCESS REQUIRED</h2><p"
      " style='color: #a1a1aa; font-family: \"JetBrains Mono\"; font-size: 12px;'>Verify your Somaiya university email to"
      " receive an access token.</p></div>",
      unsafe_allow_html=True,
  )

  col_space1, col_login, col_space3 = st.columns([3, 4, 3])
  with col_login:
    st.markdown("<div class='glass-panel'>", unsafe_allow_html=True)
    if not st.session_state.otp_sent:
      email_input = st.text_input(
          "University Email", placeholder="your.name@somaiya.edu"
      )
      if st.button("DISPATCH ACCESS CODE", use_container_width=True):
        if email_input.strip().lower().endswith("@somaiya.edu"):
          with st.spinner("Dispatching secure token..."):
            otp = str(random.randint(100000, 999999))
            st.session_state.generated_otp = otp
            st.session_state.user_email = email_input.strip().lower()
            st.session_state.otp_timestamp = time.time()
            st.session_state.otp_attempts = 0

            success, error_msg = send_otp_email(
                st.session_state.user_email, otp
            )

            if success:
              st.session_state.otp_sent = True
              st.rerun()
            else:
              st.error(f"❌ Failed to dispatch email: {error_msg}")
        else:
          st.error(
              "❌ Access Denied. Only @somaiya.edu accounts are permitted."
          )
    else:
      st.success(f"Secure token dispatched to {st.session_state.user_email}")
      otp_input = st.text_input("Enter 6-Digit Token", type="password")

      if st.button("VERIFY & ENTER", use_container_width=True):
        # 5. HARDENED AUTHENTICATION CHECKS
        curr_attempts = st.session_state.get("otp_attempts", 0)
        otp_age = time.time() - st.session_state.get("otp_timestamp", 0)

        if curr_attempts >= 5:
          st.error("❌ Too many failed attempts (5/5). Access locked. Please click 'Use a different identity' to restart.")
        elif otp_age > 600: # 10-minute OTP expiration limit
          st.error("❌ Access code has expired (valid for 10 minutes). Please click 'Use a different identity' to request a new code.")
        elif otp_input.strip() == st.session_state.generated_otp:
          st.session_state.authenticated = True
          signed_token = sign_session_token(st.session_state.user_email)
          cookie_manager.set(
              "apollo_somaiya_session",
              signed_token,
              expires_at=datetime.datetime.now()
              + datetime.timedelta(days=30),
          )
          st.rerun()
        else:
          st.session_state.otp_attempts = curr_attempts + 1
          rem = 5 - st.session_state.otp_attempts
          st.error(f"❌ Invalid token. ({rem} attempt(s) remaining)")

      st.markdown("<br>", unsafe_allow_html=True)
      if st.button("Use a different identity", type="secondary"):
        st.session_state.otp_sent = False
        st.session_state.otp_attempts = 0
        st.session_state.generated_otp = None
        st.rerun()

    st.markdown("</div>", unsafe_allow_html=True)
  st.stop()
# ==============================================================


# ================= SIDEBAR: SETTINGS & TOOLS =================
with st.sidebar:
  st.markdown(
      "<h2 class='omni-brand' style='font-size: 16px; margin-bottom: 12px;'>APOLLO <span>OMNI</span><br><span style='font-size: 9px; color: #71717a; font-family: \"JetBrains Mono\";'>SYSTEM CONTROL PANEL</span></h2>",
      unsafe_allow_html=True,
  )

  # Per-user Notebooks: switch between separate source/chat workspaces
  render_notebook_switcher(st.session_state.user_email, embedder)

  # Apply a pending cross-panel navigation jump (e.g. the Studio's Tutor tile)
  # BEFORE the radio widget below is instantiated -- setting a widget's own
  # session_state key after it has rendered this run would raise an error.
  if "nav_override" in st.session_state:
    st.session_state["main_app_navigation"] = st.session_state.pop("nav_override")

  # Interactive App Navigation Switcher
  app_mode = st.radio(
      "NAVIGATION:",
      options=[
          "⚡ Console & Tools",
          "🎓 Socratic Tutor",
          "📊 Progress Dashboard",
          "🗓️ Study Planner",
          "⚙️ User Settings & Profile",
      ],
      key="main_app_navigation",
  )
  st.markdown("<hr style='border-color: rgba(255,140,0,0.2); margin: 12px 0;'>", unsafe_allow_html=True)

  # Telemetry Row
  st.markdown(
      "<div class='glass-panel' style='padding: 12px; margin-bottom: 16px;'>",
      unsafe_allow_html=True,
  )
  c_lat, c_vec = st.columns(2)
  with c_lat:
    st.markdown(
        f"<div style='text-align:center;'><div style='font-size:20px; font-weight:900; color:white; font-family: \"JetBrains Mono\";'>{st.session_state.response_time}<span style='font-size:12px;color:#ff8c00;'>s</span></div><div style='font-size:9px; color:#a1a1aa; text-transform:uppercase;'>Latency</div></div>",
        unsafe_allow_html=True,
    )
  with c_vec:
    st.markdown(
        f"<div style='text-align:center;'><div style='font-size:20px; font-weight:900; color:white; font-family: \"JetBrains Mono\";'>{st.session_state.node_count}</div><div style='font-size:9px; color:#a1a1aa; text-transform:uppercase;'>Vectors</div></div>",
        unsafe_allow_html=True,
    )
  st.markdown("</div>", unsafe_allow_html=True)

  # Export & trust controls
  st.markdown(
      "<div class='glass-panel' style='padding: 12px; margin-bottom: 16px;'>",
      unsafe_allow_html=True,
  )
  st.markdown(
      "<div class='panel-header' style='margin-bottom: 8px; padding-bottom:"
      " 8px;'>📄 Export & Trust</div>",
      unsafe_allow_html=True,
  )
  st.session_state.answer_key_enabled = st.checkbox(
      "Generate answer-key DOCX for question papers",
      value=st.session_state.get("answer_key_enabled", True),
      help="When Apollo detects a question paper, also create a matching marking scheme document.",
  )
  st.session_state.chat_pdf_enabled = st.checkbox(
      "Show PDF exports where available",
      value=st.session_state.get("chat_pdf_enabled", True),
      help="Adds universal PDF downloads for generated papers, answers, reports, plans, and slide outlines.",
  )
  st.markdown("</div>", unsafe_allow_html=True)

  # Inference Engine Selection
  st.markdown(
      "<div class='glass-panel' style='padding: 12px; margin-bottom: 16px;'>",
      unsafe_allow_html=True,
  )
  st.markdown(
      "<div class='panel-header' style='margin-bottom: 8px; padding-bottom:"
      " 8px;'>⚙️ Inference Engine</div>",
      unsafe_allow_html=True,
  )

  model_list = list(MODEL_OPTIONS.keys())
  saved_model = st.session_state.user_prefs.get(
      "default_model", "Qwen 3.6 27B (Groq)"
  )
  default_index = (
      model_list.index(saved_model) if saved_model in model_list else 0
  )

  selected_model = st.selectbox(
      "API Gateway Endpoint:",
      options=model_list,
      index=default_index,
      label_visibility="collapsed",
  )
  st.markdown(
      f"<div style='font-size: 9px; color: #a1a1aa; font-family: \"JetBrains"
      f" Mono\"; margin-top: 4px;'>{MODEL_OPTIONS[selected_model]['desc']}</div>",
      unsafe_allow_html=True,
  )
  st.markdown("</div>", unsafe_allow_html=True)

  # Cross-notebook search
  with st.expander("🔎 Search across notebooks", expanded=False):
    cross_query = st.text_input(
        "Search everything uploaded",
        placeholder="e.g., process scheduling algorithms",
        key="cross_notebook_query",
        label_visibility="collapsed",
    )
    if st.button("Search all notebooks", use_container_width=True, key="cross_notebook_btn"):
      with st.spinner("Searching saved notebook indexes..."):
        st.session_state.cross_notebook_results = search_all_notebooks(
            st.session_state.user_email, cross_query, embedder, k=3
        )
    if st.session_state.cross_notebook_results:
      for idx, item in enumerate(st.session_state.cross_notebook_results, start=1):
        with st.expander(f"{idx}. {item['notebook']} — {item['source']}", expanded=idx == 1):
          st.write(item["snippet"])

  # Web Crawler (Tavily)
  st.markdown(
      "<div class='glass-panel' style='padding: 12px; margin-bottom: 16px;'>",
      unsafe_allow_html=True,
  )
  st.markdown(
      "<div class='panel-header' style='margin-bottom: 8px; padding-bottom:"
      " 8px;'>🌐 Web Crawler (Tavily)</div>",
      unsafe_allow_html=True,
  )
  web_query = st.text_input(
      "Query:",
      placeholder="e.g. Artificial Intelligence trends 2026",
      label_visibility="collapsed",
  )
  if st.button("FETCH & INDEX", use_container_width=True):
    if not TAVILY_API_KEY or not TAVILY_API_KEY.startswith("tvly-"):
      st.error("No active Tavily API Key found in Streamlit Secrets.")
    elif web_query:
      with st.spinner("Executing secure web retrieval..."):
        try:
          result = _cached_tavily_search(web_query, TAVILY_API_KEY.strip())

          if tavily_answer := result.get("answer"):
            st.info(f"💡 **Tavily Quick Answer:** {tavily_answer}")

          results = result.get("results", [])
          web_docs = [
              LangchainDocument(
                  page_content=f"Title: {r.get('title')}\nSource: {r.get('url')}\nContext: {r.get('content')}",
                  metadata={
                      "source": r.get("url", ""),
                      "title": r.get("title", ""),
                  },
              )
              for r in results
          ]
          chunks = text_splitter.split_documents(web_docs)
          if chunks:
            if st.session_state.vector_db is None:
              st.session_state.vector_db = FAISS.from_documents(
                  chunks, embedder
              )
            else:
              st.session_state.vector_db.add_documents(chunks)
            st.session_state.node_count += len(chunks)
            for r in results:
              register_source(r.get("url") or r.get("title") or "Web result", kind="web")
            st.success(f"Indexed {len(chunks)} blocks!")
            save_active_notebook(st.session_state.user_email)
          elif not results:
            st.warning("⚠️ No results returned. Try a different query.")
        except Exception as e:
          st.error(f"Search failed: {str(e)}")
  st.markdown("</div>", unsafe_allow_html=True)

  # Local Documents Upload
  st.markdown(
      "<div class='glass-panel' style='padding: 12px; margin-bottom: 16px;'>",
      unsafe_allow_html=True,
  )
  st.markdown(
      "<div class='panel-header' style='margin-bottom: 8px; padding-bottom:"
      " 8px;'>📚 Material Append</div>",
      unsafe_allow_html=True,
  )
  uploaded_files = st.file_uploader(
      "Upload course materials...",
      type=["pdf", "txt", "docx"],
      accept_multiple_files=True,
      label_visibility="collapsed",
      key="file_in",
  )
  force_reindex = False
  if uploaded_files:
    _already = set(get_source_names())
    _dupes = [f.name for f in uploaded_files if f.name in _already]
    if _dupes:
      force_reindex = st.checkbox(
          f"🔁 Re-index {len(_dupes)} file(s) already in this notebook",
          value=False,
          help="Off = skip files with a name that's already indexed here.",
      )

  if st.button("+ ADD TO KNOWLEDGE", use_container_width=True):
    if not uploaded_files:
      st.warning("Choose at least one PDF, TXT, or DOCX file first.")
    else:
      already_indexed = set(get_source_names())
      to_parse = [
          f for f in uploaded_files
          if force_reindex or f.name not in already_indexed
      ]
      skipped = [f.name for f in uploaded_files if f not in to_parse]

      docs = []
      failures = []
      if to_parse:
        progress = st.progress(0.0, text=f"Parsing 0/{len(to_parse)} files...")
        # Parse files concurrently -- I/O + PDF-parsing bound, so threads
        # give a real wall-clock speedup on multi-file uploads.
        with ThreadPoolExecutor(max_workers=min(8, len(to_parse))) as pool:
          futures = {
              pool.submit(_parse_uploaded_file, f.name, f.read()): f.name
              for f in to_parse
          }
          done = 0
          for future in as_completed(futures):
            f_name, loaded, err = future.result()
            done += 1
            progress.progress(done / len(to_parse), text=f"Parsing {done}/{len(to_parse)} files...")
            if err:
              failures.append(f"{f_name} — {err}")
            else:
              docs.extend(loaded)
              register_source(f_name, kind="file")
        progress.empty()

      if docs:
        chunks = text_splitter.split_documents(docs)
        if chunks:
          # Embed in batches with a progress bar instead of one giant blocking
          # call -- keeps memory use flat and gives real feedback on large
          # uploads instead of a frozen spinner.
          BATCH = 64
          embed_progress = st.progress(0.0, text=f"Embedding 0/{len(chunks)} blocks...")
          for start in range(0, len(chunks), BATCH):
            batch = chunks[start:start + BATCH]
            if st.session_state.vector_db is None:
              st.session_state.vector_db = FAISS.from_documents(batch, embedder)
            else:
              st.session_state.vector_db.add_documents(batch)
            done_n = min(start + BATCH, len(chunks))
            embed_progress.progress(done_n / len(chunks), text=f"Embedding {done_n}/{len(chunks)} blocks...")
          embed_progress.empty()

          st.session_state.node_count += len(chunks)
          st.success(f"✅ Indexed {len(chunks)} blocks from {len(to_parse) - len(failures)} file(s).")
          save_active_notebook(st.session_state.user_email)

      if skipped:
        st.caption(f"⏭️ Skipped (already indexed): {', '.join(skipped)}")
      if failures:
        st.warning("⚠️ Couldn't index:\n" + "\n".join(f"- {f}" for f in failures))
  st.markdown("</div>", unsafe_allow_html=True)

  # Session Control
  st.markdown(
      "<div class='glass-panel' style='padding: 12px;'>", unsafe_allow_html=True
  )
  st.markdown(
      "<div class='panel-header' style='margin-bottom: 8px; padding-bottom:"
      " 8px;'>🛠️ Session Control</div>",
      unsafe_allow_html=True,
  )
  if st.button("FLUSH MEMORY", use_container_width=True):
    st.session_state.chat_history = []
    st.session_state.vector_db = None
    st.session_state.node_count = 0
    st.session_state.response_time = "0.00"
    st.session_state.indexed_sources = []
    st.session_state.studio_results = {}
    st.session_state.dialog_open = False
    st.session_state.source_reference = (
        "<div class='source-box font-mono'>Awaiting vector alignment...</div>"
    )
    save_active_notebook(st.session_state.user_email)
    st.rerun()

  st.session_state.voice_output_enabled = st.checkbox(
      "🔊 Voice Output (TTS)",
      value=st.session_state.get("voice_output_enabled", False),
      help="Speak assistant responses aloud using Microsoft Edge TTS.",
  )

  if st.button("TERMINATE SESSION", use_container_width=True, type="secondary"):
    st.session_state.authenticated = False
    cookie_manager.delete("apollo_somaiya_session")
    st.rerun()
  st.markdown("</div>", unsafe_allow_html=True)


# ================= MAIN AREA: EXPANDED CONSOLE & NOTEBOOKLM STUDIO =================
if app_mode == "⚙️ User Settings & Profile":
  render_settings_page()
elif app_mode == "🎓 Socratic Tutor":
  render_tutor_mode(
      groq_key=GROQ_API_KEY,
      selected_model=selected_model,
      generate_llm_response_fn=functools.partial(generate_llm_response, gemini_key=GEMINI_API_KEY),
      generate_llm_stream_fn=functools.partial(generate_llm_stream, gemini_key=GEMINI_API_KEY),
      get_scoped_context_fn=get_scoped_context,
      source_names=get_source_names(),
      user_prefs=st.session_state.get("user_prefs"),
  )
elif app_mode == "📊 Progress Dashboard":
  render_progress_dashboard(st.session_state.user_email)
elif app_mode == "🗓️ Study Planner":
  render_study_planner(st.session_state.user_email)
else:
  # FIX 1: Column definition and ALL console/studio rendering blocks are inside else:
  col_chat, col_tools = st.columns([72, 28], gap="medium")

  # ----------------- MAIN LEFT: EXPANDED CHAT CONSOLE -----------------
  with col_chat:

    _hdr_left, _hdr_right = st.columns([5, 1])
    with _hdr_left:
      _active_nb = get_active_notebook(st.session_state.user_email)
      _nb_label = f" · {_active_nb['title']}" if _active_nb else ""
      st.markdown(
          f"""
      <div style='background: rgba(0,0,0,0.6); padding: 12px 20px; border-bottom: 1px solid rgba(255,255,255,0.05); border-radius: 6px 6px 0 0; display: flex; justify-content: space-between; align-items: center;'>
          <div style='display: flex; gap: 8px; align-items: center;'>
              <div style='width: 8px; height: 8px; background: #ef4444; border-radius: 50%; opacity: 0.8;'></div>
              <div style='width: 8px; height: 8px; background: #ff8c00; border-radius: 50%; opacity: 0.8;'></div>
              <div style='width: 8px; height: 8px; background: #22c55e; border-radius: 50%; opacity: 0.8;'></div>
              <span style='font-size: 11px; font-weight: 700; letter-spacing: 0.2em; color: #a1a1aa; text-transform: uppercase; margin-left: 12px;'>STUDY_CONSOLE_EXPANDED{_nb_label}</span>
          </div>
      </div>
      """,
          unsafe_allow_html=True,
      )

    with _hdr_right:
      if st.session_state.chat_history:
        _md_lines = []
        for _m in st.session_state.chat_history:
          _role_label = "**You**" if _m["role"] == "user" else "**Apollo**"
          _md_lines.append(f"{_role_label}:\n{_m['content']}\n")
        _md_export = "\n---\n".join(_md_lines)
        st.download_button(
            label="⬇ Export",
            data=_md_export,
            file_name="apollo_chat_transcript.md",
            mime="text/markdown",
            use_container_width=True,
            help="Download this conversation as a Markdown transcript",
        )

    if not st.session_state.chat_history:
      st.markdown(
          """
          <div style='margin-top: 70px; margin-bottom: 40px; text-align: center;'>
              <h2 style='color: #ff8c00; font-family: "Inter", sans-serif; font-weight: 700; font-size: 26px; letter-spacing: 0.1em;'>STUDY CONSOLE READY</h2>
              <p style='color: #a1a1aa; font-family: "JetBrains Mono", monospace; font-size: 13px; margin-top: 10px;'>Ask questions, analyze uploaded materials, or generate NotebookLM overview assets on the right.</p>
          </div>
          """,
          unsafe_allow_html=True,
      )

    # INCREASED CHAT BOX AREA: Height expanded to 620px
    chat_scroll_pane = st.container(height=620, border=False)

    with chat_scroll_pane:
      for _msg_idx, msg in enumerate(st.session_state.chat_history):
        with st.chat_message(msg["role"]):
          if msg["role"] == "assistant":
            _citation_sources = msg.get("citation_sources", [])
            _renderable_content = decorate_citation_links(msg["content"], _citation_sources)
            render_dynamic_chart_from_text(_renderable_content)
            render_citation_source_panel(_citation_sources)
            _docx_path = msg.get("docx_path")
            if _docx_path and os.path.exists(_docx_path):
              with open(_docx_path, "rb") as _qp_f:
                st.download_button(
                    "📄 Download as Word Doc (with space to solve)",
                    _qp_f.read(),
                    file_name="apollo_question_paper.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key=f"dl_qp_{_msg_idx}",
                )
            _answer_key_path = msg.get("answer_key_path")
            if _answer_key_path and os.path.exists(_answer_key_path):
              with open(_answer_key_path, "rb") as _ak_f:
                st.download_button(
                    "✅ Download answer key / marking scheme",
                    _ak_f.read(),
                    file_name="apollo_answer_key.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True,
                    key=f"dl_ak_{_msg_idx}",
                )
            _pdf_path = msg.get("pdf_path")
            if _pdf_path and os.path.exists(_pdf_path):
              with open(_pdf_path, "rb") as _pdf_f:
                st.download_button(
                    "📕 Download as PDF",
                    _pdf_f.read(),
                    file_name="apollo_export.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key=f"dl_pdf_{_msg_idx}",
                )
          else:
            st.markdown(msg["content"])

    # --- ASK WITH A PHOTO (handwriting / diagrams / textbook pages) ---
    render_image_question_widget(
        GROQ_API_KEY,
        key_suffix="chat_main",
        on_answered=lambda: save_active_notebook(st.session_state.user_email),
    )
    render_homework_grading_widget(
        GROQ_API_KEY,
        key_suffix="chat_main",
        on_answered=lambda: save_active_notebook(st.session_state.user_email),
    )

    # --- VOICE & TEXT INPUT MATRIX ---
    voice_prompt = render_voice_input(GROQ_API_KEY, key_suffix="chat_main")
    user_query = st.chat_input("AWAITING COMMAND OR QUESTION...")

    final_query = voice_prompt if voice_prompt else user_query

    if final_query:
      st.session_state.chat_history.append(
          {"role": "user", "content": final_query}
      )
      start_time = time.time()
      context_payload = ""
      citation_sources = []

      chart_instruction = (
          "\n\nIf the user asks for a chart, graph, data visualization, or"
          " numerical comparison, append a JSON code block at the very end of"
          ' your response following this exact structure:\n```json\n{\n  "type":'
          ' "bar",  // options: "bar", "line", or "pie"\n  "title": "Chart'
          ' Title",\n  "x_label": "X Axis Label",\n  "y_label": "Y Axis'
          ' Label",\n  "x": ["Category A", "Category B"],\n  "y": [10, 20]\n}\n```'
      )

      _prefs = st.session_state.get("user_prefs", {})
      _style = _prefs.get("learning_style", "General")
      _depth = _prefs.get("detail_level", "Intermediate")
      _name  = _prefs.get("full_name", "").strip()
      prefs_preamble = (
          f"Student profile: learning style = '{_style}', "
          f"detail level = '{_depth}'."
          + (f" Address the student as {_name}." if _name else "")
          + " Tailor all responses accordingly.\n\n"
      )

      _auto_search_fired = False
      if st.session_state.vector_db is None and _needs_web_search(final_query):
        if TAVILY_API_KEY and TAVILY_API_KEY.startswith("tvly-"):
          try:
            with st.spinner("🌐 Fetching real-time context via Tavily..."):
              _auto_result = _cached_tavily_search(final_query, TAVILY_API_KEY.strip())

            if _ta := _auto_result.get("answer"):
              st.info(f"💡 **Tavily Quick Answer:** {_ta}")

            _web_ctx_parts = [
                f"[{r.get('title', 'Web Result')}]\nSource: {r.get('url', '')}\n{r.get('content', '')}"
                for r in _auto_result.get("results", [])[:3]
            ]
            if _web_ctx_parts:
              context_payload = "\n\n".join(_web_ctx_parts)
              _auto_search_fired = True

              # Build lightweight citation sources from Tavily results so the UI
              # can render Source panels and inline [source: S1] tags.
              citation_sources = []
              for _idx, _r in enumerate(_auto_result.get("results", [])[:3], start=1):
                _snippet = re.sub(r"\s+", " ", _r.get("content", "")).strip()
                citation_sources.append({
                  "id": f"S{_idx}",
                  "label": _r.get("title", "Web Result"),
                  "source": _r.get("url", ""),
                  "page": None,
                  "snippet": _snippet[:1200],
                })

              _clean_web_ctx = (
                  context_payload.replace("<", "&lt;")
                  .replace(">", "&gt;")
                  .replace("\n", "<br>")
              )
              st.session_state.source_reference = (
                  "<div class='source-box'><strong>Active Context"
                  " (Tavily Web Search):</strong><br><br>"
                  f"{_clean_web_ctx}</div>"
              )
          except Exception as _ae:
            st.warning(f"⚠️ Auto-search failed gracefully: {_ae}")

      if st.session_state.vector_db is not None:
        retriever = st.session_state.vector_db.as_retriever(
            search_kwargs={"k": 5}
        )
        matched_nodes = retriever.invoke(final_query)
        citation_sources = build_citation_sources(matched_nodes)
        context_payload = citation_context_from_sources(citation_sources)
        sys_instruction = (
            f"{prefs_preamble}You are APOLLO OMNI AI, an advanced study assistant powered by"
            f" Groq LPUs. Answer using ONLY context below. Cite important claims inline"
            f" with the exact tag for the chunk you used, such as [source: S1]."
            f" Use only source IDs present in the Context Matrix.{chart_instruction}"
        )
        clean_ctx = (
            context_payload.replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace("\n", "<br>")
        )
        st.session_state.source_reference = (
            "<div class='source-box'><strong>Active Context"
            f" (RAG):</strong><br><br>{clean_ctx}</div>"
        )
      elif _auto_search_fired:
        sys_instruction = (
            f"{prefs_preamble}You are APOLLO OMNI AI, an advanced study assistant powered by"
            f" Groq LPUs. Use the real-time web context below to answer accurately.{chart_instruction}"
        )
      else:
        sys_instruction = (
            f"{prefs_preamble}You are APOLLO OMNI AI, an advanced study assistant powered by"
            f" Groq LPUs. Answer based on general knowledge.{chart_instruction}"
        )
        st.session_state.source_reference = (
            "<div class='source-box font-mono'>No active context. General weights"
            " used.</div>"
        )

      message_stream = [{"role": "system", "content": sys_instruction}]
      for msg in st.session_state.chat_history[-4:]:
        message_stream.append({"role": msg["role"], "content": msg["content"]})
      message_stream.append({
          "role": "user",
          "content": f"Context Matrix:\n{context_payload}\n\nQuery: {final_query}",
      })

      with chat_scroll_pane:
        with st.chat_message("assistant"):
          try:
            stream = generate_llm_stream(
                message_stream,
                GROQ_API_KEY,
                selected_model,
                gemini_key=GEMINI_API_KEY,
            )
            collected_tokens = st.write_stream(stream)
            if not collected_tokens or not str(collected_tokens).strip():
              collected_tokens = "⚠️ EMPTY RESPONSE."
              st.markdown(collected_tokens)
          except Exception as ex:
            collected_tokens = f"❌ FRAMEWORK API FAILURE: {ex}"
            st.markdown(collected_tokens)

          # FIX 6: Standardized TTS call with explicit voice parameter
          if st.session_state.get("voice_output_enabled", False):
            _tts_text = str(collected_tokens).strip()
            if _tts_text and not _tts_text.startswith("❌"):
              with st.spinner("🎶 Synthesizing voice..."):
                _audio_bytes = run_tts_synthesis(_tts_text, voice="en-US-AriaNeural")
              if _audio_bytes:
                st.audio(_audio_bytes, format="audio/mp3")

      qp_docx_path = None
      answer_key_path = None
      pdf_path = None
      is_question_paper = _looks_like_question_paper(final_query)
      if is_question_paper and collected_tokens and not str(collected_tokens).startswith("❌"):
        try:
          qp_docx_path = docx_from_chat_answer(final_query[:120], str(collected_tokens))
        except Exception:
          qp_docx_path = None
        if st.session_state.get("answer_key_enabled", True):
          try:
            key_prompt = (
                "Create a concise answer key and marking scheme for the question paper below. "
                "Match the question numbering, include expected answer points, and add suggested marks "
                "where the paper implies them.\n\n"
                f"QUESTION PAPER:\n{collected_tokens}"
            )
            answer_key_md, _key_status = generate_llm_response(
                [{"role": "user", "content": key_prompt}],
                GROQ_API_KEY,
                selected_model,
                max_tokens=1600,
                gemini_key=GEMINI_API_KEY,
            )
            if answer_key_md:
              answer_key_path = answer_key_docx_from_question_paper(final_query[:120], answer_key_md)
          except Exception:
            answer_key_path = None

      if st.session_state.get("chat_pdf_enabled", True) and collected_tokens and not str(collected_tokens).startswith("❌"):
        try:
          pdf_path = pdf_from_markdown(final_query[:120], str(collected_tokens), ruled_questions=is_question_paper)
        except Exception:
          pdf_path = None

      st.session_state.chat_history.append({
          "role": "assistant",
          "content": collected_tokens,
          "docx_path": qp_docx_path,
          "answer_key_path": answer_key_path,
          "pdf_path": pdf_path,
          "citation_sources": citation_sources,
      })
      st.session_state.response_time = f"{time.time() - start_time:.2f}"
      save_active_notebook(st.session_state.user_email)
      st.rerun()


  # ----------------- MAIN RIGHT: NOTEBOOKLM STYLE STUDIO GRID -----------------
  with col_tools:

    st.markdown("<div class='glass-panel'>", unsafe_allow_html=True)
    st.markdown(
        "<div class='panel-header' style='font-size: 13px; letter-spacing: 0.25em;'>⚡ STUDIO (NOTEBOOK LM OVERVIEW)</div>",
        unsafe_allow_html=True,
    )

    # Banner prompt inside Studio header
    st.markdown(
        """
        <div style='background: linear-gradient(135deg, rgba(255, 140, 0, 0.15), rgba(30, 30, 40, 0.8)); border: 1px solid rgba(255, 140, 0, 0.3); border-radius: 8px; padding: 12px 16px; margin-bottom: 16px;'>
            <div style='font-size: 11px; font-weight: 700; color: #ff8c00; font-family: "JetBrains Mono", monospace;'>
                ✨ GENERATE STUDIO OVERVIEW
            </div>
            <div style='font-size: 10px; color: #a1a1aa; margin-top: 4px;'>
                Select any studio feature tile below to open its creation window, scoped to your indexed sources.
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    # ── NOTEBOOK LM 2-COLUMN FEATURE GRID ─────────────────────────────────────
    g_col1, g_col2 = st.columns(2)

    with g_col1:
      btn_audio = st.button(
          "🎙️ Audio Overview  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Audio Overview" else "primary",
          key="tile_audio",
      )
      btn_video = st.button(
          "🎬 Video Overview  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Video Overview" else "primary",
          key="tile_video",
      )
      btn_reports = st.button(
          "📝 Study Reports  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Study Reports" else "primary",
          key="tile_reports",
      )
      btn_quiz = st.button(
          "❓ Practice Quiz  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Practice Quiz" else "primary",
          key="tile_quiz",
      )

    with g_col2:
      btn_slides = st.button(
          "💻 Slide Deck [BETA] ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Slide Deck" else "primary",
          key="tile_slides",
      )
      btn_mindmap = st.button(
          "🧠 Mind Map  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Mind Map" else "primary",
          key="tile_mindmap",
      )
      btn_flashcards = st.button(
          "📇 Flashcards  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Flashcards" else "primary",
          key="tile_flashcards",
      )
      btn_context = st.button(
          "📑 Active Context  ›",
          use_container_width=True,
          type="secondary" if st.session_state.active_studio_tool != "Active Context" else "primary",
          key="tile_context",
      )

    # Clicking a generation tile opens its NotebookLM-style creation dialog;
    # "Active Context" just switches the persistent panel below (no dialog needed).
    if btn_audio:
      st.session_state.active_studio_tool = "Audio Overview"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_slides:
      st.session_state.active_studio_tool = "Slide Deck"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_video:
      st.session_state.active_studio_tool = "Video Overview"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_mindmap:
      st.session_state.active_studio_tool = "Mind Map"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_reports:
      st.session_state.active_studio_tool = "Study Reports"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_flashcards:
      st.session_state.active_studio_tool = "Flashcards"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_quiz:
      st.session_state.active_studio_tool = "Practice Quiz"
      st.session_state.dialog_open = True
      st.rerun()
    elif btn_context:
      st.session_state.active_studio_tool = "Active Context"
      st.session_state.dialog_open = False
      st.rerun()

    st.markdown(
        "<div style='font-size:9px; color:#71717a; text-align:center; margin: 6px 0 4px 0;"
        " text-transform: uppercase; letter-spacing: 0.1em;'>Beyond one-shot generation</div>",
        unsafe_allow_html=True,
    )
    if st.button("🎓  SOCRATIC TUTOR — Get Tested & Taught Live  ›", use_container_width=True, key="tile_tutor", type="primary"):
      st.session_state["nav_override"] = "🎓 Socratic Tutor"
      st.rerun()

    st.markdown("<hr style='border-color: rgba(255,140,0,0.2); margin: 16px 0;'>", unsafe_allow_html=True)

    # ── OPEN THE NOTEBOOKLM-STYLE CREATION DIALOG FOR THE ACTIVE TOOL ───────
    active_tool = st.session_state.get("active_studio_tool", "Slide Deck")
    if st.session_state.dialog_open and active_tool in _STUDIO_DIALOGS:
      _STUDIO_DIALOGS[active_tool]()

    # ── PERSISTENT RESULTS PANEL (renders whatever was last generated) ──────
    result = st.session_state.studio_results.get(active_tool)

    # 1. SLIDE DECK (Gamma-style presentation generator)
    if active_tool == "Slide Deck":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>💻 PRESENTATION SLIDE DECK</div>",
          unsafe_allow_html=True,
      )
      if not st.session_state.slides_data or not isinstance(st.session_state.slides_data, list):
        st.session_state.slides_data = [{
            "title": "Welcome to Apollo Omni AI",
            "subtitle": "Awaiting Presentation Prompt",
            "image_keyword": "abstract technology minimalist",
            "image_prompt": (
                "abstract technology network, dark navy background, orange"
                " neon highlights, cinematic lighting, no text"
            ),
            "cards": [{"heading": "Getting Started", "text": "Click 'Slide Deck' above to generate a deck from your sources."}],
        }]

      if st.button("💻 Open Slide Deck Generator", use_container_width=True, key="reopen_slides"):
        st.session_state.dialog_open = True
        st.rerun()

      with st.expander("✏️ Live Slide Editor", expanded=True):
        tabs = st.tabs([f"S{i+1}" for i in range(len(st.session_state.slides_data))])
        for i, tab in enumerate(tabs):
          with tab:
            slide_info = st.session_state.slides_data[i]
            st.session_state.slides_data[i]["title"] = st.text_input(
                f"Title {i+1}", slide_info.get("title", ""), key=f"t_{i}"
            )
            st.session_state.slides_data[i]["subtitle"] = st.text_input(
                f"Subtitle {i+1}", slide_info.get("subtitle", ""), key=f"sub_{i}"
            )
            current_img = slide_info.get("image_prompt") or slide_info.get(
                "image_keyword", ""
            )
            updated_img = st.text_input(
                f"Image prompt {i+1}", current_img, key=f"img_{i}"
            )
            st.session_state.slides_data[i]["image_prompt"] = updated_img
            st.session_state.slides_data[i]["image_keyword"] = updated_img

            cards = slide_info.get("cards", [])
            if not isinstance(cards, list):
              cards = [{"heading": "Detail", "text": str(cards)}]

            for j, card in enumerate(cards):
              st.markdown(
                  f"<div style='font-size: 10px; font-weight: bold; margin-top: 8px; color: #a1a1aa;'>Card {j+1}</div>",
                  unsafe_allow_html=True,
              )
              if isinstance(card, dict):
                cards[j]["heading"] = st.text_input(
                    f"Heading", card.get("heading", ""), key=f"ch_{i}_{j}", label_visibility="collapsed"
                )
                cards[j]["text"] = st.text_area(
                    f"Text", card.get("text", ""), key=f"ct_{i}_{j}", label_visibility="collapsed"
                )
            st.session_state.slides_data[i]["cards"] = cards

      if st.button("📥 EXPORT .PPTX DECK", use_container_width=True):
        with st.spinner("Building PowerPoint file..."):
          file_path = create_gamma_style_pptx(st.session_state.slides_data)
          with open(file_path, "rb") as f:
            st.download_button(
                label="DOWNLOAD FILE",
                data=f,
                file_name="Apollo_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
      if st.session_state.get("chat_pdf_enabled", True) and st.button("📕 EXPORT .PDF OUTLINE", use_container_width=True):
        pdf_path = pdf_from_slides_data(st.session_state.slides_data)
        if pdf_path and os.path.exists(pdf_path):
          with open(pdf_path, "rb") as f:
            st.download_button(
                label="DOWNLOAD PDF",
                data=f.read(),
                file_name="Apollo_Presentation_Outline.pdf",
                mime="application/pdf",
                use_container_width=True,
            )
        else:
          st.warning("PDF export needs the reportlab package installed.")

    # 2. VIDEO OVERVIEW
    elif active_tool == "Video Overview":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>🎬 AI VIDEO OVERVIEW GENERATOR</div>",
          unsafe_allow_html=True,
      )
      st.caption("Click 'Video Overview' above to open the generator window.")
      if st.button("🎬 Open Video Generator", use_container_width=True, key="reopen_video"):
        st.session_state.dialog_open = True
        st.rerun()

    # 3. AUDIO OVERVIEW (NotebookLM Podcast Style Audio Summary)
    elif active_tool == "Audio Overview":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>🎙️ AUDIO OVERVIEW (PODCAST SYNTHESIS)</div>",
          unsafe_allow_html=True,
      )
      if st.button("🎙️ Open Audio Overview Generator", use_container_width=True, key="reopen_audio"):
        st.session_state.dialog_open = True
        st.rerun()

      if result:
        if result.get("sources"):
          st.caption(f"Based on: {', '.join(result['sources'])}")
        st.markdown(f"**Generated Script:**\n\n{result.get('script', '')}")
        if result.get("audio"):
          st.audio(result["audio"], format="audio/mp3")
          st.download_button(
              "📥 DOWNLOAD AUDIO (MP3)", result["audio"],
              file_name="notebooklm_audio_overview.mp3", mime="audio/mp3",
              use_container_width=True, key="dl_audio_overview",
          )
      else:
        st.info("Nothing generated yet — open the generator above to create an audio overview.")

    # 4. MIND MAP GENERATOR
    elif active_tool == "Mind Map":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>🧠 CONCEPT MIND MAP GENERATOR</div>",
          unsafe_allow_html=True,
      )
      if st.button("🧠 Open Mind Map Generator", use_container_width=True, key="reopen_mindmap"):
        st.session_state.dialog_open = True
        st.rerun()

      if result:
        if result.get("sources"):
          st.caption(f"Based on: {', '.join(result['sources'])}")
        st.code(result.get("content", ""), language="markdown")
      else:
        st.info("Nothing generated yet — open the generator above to create a mind map.")

    # 5. STUDY REPORTS
    elif active_tool == "Study Reports":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>📝 COMPREHENSIVE STUDY REPORT</div>",
          unsafe_allow_html=True,
      )
      if st.button("📝 Open Study Report Generator", use_container_width=True, key="reopen_reports"):
        st.session_state.dialog_open = True
        st.rerun()

      if result:
        if result.get("sources"):
          st.caption(f"Based on: {', '.join(result['sources'])}")
        st.markdown(result.get("content", ""))
        st.download_button(
            "📥 DOWNLOAD REPORT (.MD)", result.get("content", ""),
            file_name="apollo_study_report.md", mime="text/markdown",
            use_container_width=True, key="dl_report",
        )
        if st.session_state.get("chat_pdf_enabled", True):
          pdf_path = pdf_from_markdown(result.get("topic", "Apollo Study Report"), result.get("content", ""))
          if pdf_path and os.path.exists(pdf_path):
            with open(pdf_path, "rb") as f:
              st.download_button(
                  "📕 DOWNLOAD REPORT (.PDF)",
                  f.read(),
                  file_name="apollo_study_report.pdf",
                  mime="application/pdf",
                  use_container_width=True,
                  key="dl_report_pdf",
              )
      else:
        st.info("Nothing generated yet — open the generator above to create a study report.")

    # 6. FLASHCARDS
    elif active_tool == "Flashcards":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>📇 REVISION FLASHCARDS</div>",
          unsafe_allow_html=True,
      )
      if st.button("📇 Open Flashcard Generator", use_container_width=True, key="reopen_flashcards"):
        st.session_state.dialog_open = True
        st.rerun()

      if result:
        if result.get("sources"):
          st.caption(f"Based on: {', '.join(result['sources'])}")
        if result.get("cards_added"):
          st.success(f"Saved {result['cards_added']} card(s) into spaced repetition.")
        st.markdown(result.get("content", ""))
        if st.session_state.get("chat_pdf_enabled", True):
          pdf_path = pdf_from_markdown(result.get("topic", "Apollo Flashcards"), result.get("content", ""))
          if pdf_path and os.path.exists(pdf_path):
            with open(pdf_path, "rb") as f:
              st.download_button(
                  "📕 DOWNLOAD FLASHCARDS (.PDF)",
                  f.read(),
                  file_name="apollo_flashcards.pdf",
                  mime="application/pdf",
                  use_container_width=True,
                  key="dl_flashcards_pdf",
              )
      else:
        st.info("Nothing generated yet — open the generator above to create flashcards.")

      st.markdown("<hr style='border-color: rgba(255,140,0,0.15); margin: 12px 0;'>", unsafe_allow_html=True)
      due_cards = due_flashcards(st.session_state.user_email)
      all_cards = get_user_flashcards(st.session_state.user_email)
      st.markdown(f"**Review Queue:** {len(due_cards)} due today · {len(all_cards)} saved")
      if due_cards:
        card = due_cards[0]
        with st.expander(f"Due now: {card.get('topic', 'General')}", expanded=True):
          st.markdown(f"**Q:** {card.get('question', '')}")
          with st.expander("Reveal answer", expanded=False):
            st.markdown(card.get("answer", ""))
          q1, q2, q3, q4 = st.columns(4)
          grades = [(q1, "Again", 1), (q2, "Hard", 3), (q3, "Good", 4), (q4, "Easy", 5)]
          for col, label, quality in grades:
            with col:
              if st.button(label, use_container_width=True, key=f"review_{card['id']}_{quality}"):
                review_flashcard(st.session_state.user_email, card["id"], quality)
                st.rerun()
      elif all_cards:
        next_due = min(c.get("due", datetime.date.today().isoformat()) for c in all_cards)
        st.caption(f"No cards due. Next review: {next_due}.")

    # 7. PRACTICE QUIZ
    elif active_tool == "Practice Quiz":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>❓ PRACTICE QUIZ GENERATOR</div>",
          unsafe_allow_html=True,
      )
      if st.button("❓ Open Quiz Generator", use_container_width=True, key="reopen_quiz"):
        st.session_state.dialog_open = True
        st.rerun()

      if result:
        if result.get("sources"):
          st.caption(f"Based on: {', '.join(result['sources'])}")
        st.markdown(result.get("content", ""))
        if st.session_state.get("chat_pdf_enabled", True):
          pdf_path = pdf_from_markdown(result.get("topic", "Apollo Practice Quiz"), result.get("content", ""))
          if pdf_path and os.path.exists(pdf_path):
            with open(pdf_path, "rb") as f:
              st.download_button(
                  "📕 DOWNLOAD QUIZ (.PDF)",
                  f.read(),
                  file_name="apollo_practice_quiz.pdf",
                  mime="application/pdf",
                  use_container_width=True,
                  key="dl_quiz_pdf",
              )
      else:
        st.info("Nothing generated yet — open the generator above to create a practice quiz.")

    # 8. ACTIVE CONTEXT
    elif active_tool == "Active Context":
      st.markdown(
          "<div style='font-size: 12px; font-weight: 700; color: #ff8c00; font-family: \"JetBrains Mono\", monospace; margin-bottom: 8px;'>📑 ACTIVE VECTOR RAG CONTEXT</div>",
          unsafe_allow_html=True,
      )
      if st.session_state.indexed_sources:
        st.markdown(
            "<div style='font-size: 10px; color: #a1a1aa; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 6px;'>Indexed Sources</div>",
            unsafe_allow_html=True,
        )
        for s in st.session_state.indexed_sources:
          icon = "📄" if s["kind"] == "file" else "🌐"
          st.markdown(f"{icon} {s['name']}")
        st.markdown("<hr style='border-color: rgba(255,140,0,0.15); margin: 10px 0;'>", unsafe_allow_html=True)
      st.markdown(st.session_state.source_reference, unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)


